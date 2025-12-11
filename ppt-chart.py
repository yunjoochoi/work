import pandas as pd
from pptx import Presentation

def get_chart_context(shape, all_shapes):
    """
    차트(shape)를 기준으로 위/아래 가장 가까운 텍스트를 찾습니다.
    """
    chart_top = shape.top
    chart_bottom = shape.top + shape.height
    
    text_above = []
    text_below = []

    for other in all_shapes:
        # 자기 자신(차트)이거나 텍스트가 없으면 패스
        if other == shape or not other.has_text_frame:
            continue
            
        text_content = other.text_frame.text.strip()
        if not text_content:
            continue

        other_top = other.top
        other_bottom = other.top + other.height

        # 1. 차트보다 위에 있는 텍스트 (텍스트의 밑변 < 차트의 윗변)
        if other_bottom < chart_top:
            distance = chart_top - other_bottom
            text_above.append((distance, text_content))
            
        # 2. 차트보다 아래에 있는 텍스트 (텍스트의 윗변 > 차트의 아랫변)
        elif other_top > chart_bottom:
            distance = other_top - chart_bottom
            text_below.append((distance, text_content))

    # 거리순 정렬 (가장 가까운 것이 0번 인덱스)
    text_above.sort(key=lambda x: x[0])
    text_below.sort(key=lambda x: x[0])

    # 가장 가까운 텍스트 반환 (없으면 None)
    context_pre = text_above[0][1] if text_above else ""
    context_post = text_below[0][1] if text_below else ""
    
    return context_pre, context_post

def extract_chart_with_context(pptx_path):
    prs = Presentation(pptx_path)
    results = []

    for slide_idx, slide in enumerate(prs.slides):
        # 해당 슬라이드의 모든 도형을 리스트로 확보
        all_shapes = list(slide.shapes)
        
        for shape in slide.shapes:
            if shape.has_chart:
                chart = shape.chart
                try:
                    title = chart.chart_title.text_frame.text
                except:
                    title = "제목 없음"

                # --- [핵심] 앞뒤 문맥 찾기 ---
                pre_text, post_text = get_chart_context(shape, all_shapes)

                # --- 차트 데이터 마크다운 변환 (간소화) ---
                try:
                    df = pd.DataFrame()
                    plot = chart.plots[0]
                    cats = [c.label for c in plot.categories]
                    # 카테고리 없으면 임의 생성
                    if not cats: cats = [f"Item {i}" for i in range(len(plot.series[0].values))]
                        
                    df.index = cats
                    for ser in plot.series:
                        df[ser.name] = pd.Series(ser.values, index=cats)
                    md_table = df.to_markdown()
                except:
                    md_table = "(데이터 추출 실패)"

                # --- 최종 RAG용 청크 생성 ---
                chunk = f"""
### 슬라이드 {slide_idx+1} 분석
**[관련 설명 (위)]**
{pre_text}

**[차트: {title}]**
{md_table}

**[관련 설명 (아래)]**
{post_text}
--------------------------------------------------
"""
                results.append(chunk)
                print(f"✅ 슬라이드 {slide_idx+1} 처리 완료")

    return "\n".join(results)

# 실행
file_path = "yr_path"
print(extract_chart_with_context(file_path))