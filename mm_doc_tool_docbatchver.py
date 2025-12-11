import shutil
import time
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, Optional, List, Tuple

from docling.document_converter import DocumentConverter, PdfFormatOption
from docling.datamodel.base_models import InputFormat
from docling.datamodel.pipeline_options import PdfPipelineOptions
from docling_core.types.doc.base import ImageRefMode
from docling_core.types.io import DocumentStream
from docling.datamodel.settings import settings

from dataclasses import dataclass

# for _process_pdf_document
# from docling_core.transforms.serializer.markdown import MarkdownDocSerializer
# from docling_core.types.doc import DocItem, PictureItem, TableItem
# import fitz  
# from markitdown import MarkItDown

from openpyxl import load_workbook
from openpyxl.utils import range_boundaries
from pptx import Presentation
import pandas as pd
import re



@dataclass
class ParserConfig:
    """
    Configuration dataclass for DoclingParser.

    Attributes:
        regions_dir: Base directory path for storing extracted images
        do_ocr: Whether to perform OCR on images within documents
        do_table_structure: Whether to detect and preserve table structures
        images_scale: Scale factor for image resolution (higher = better quality)
        layout_batch_size: Batch size for layout detection processing
        table_batch_size: Batch size for table structure recognition
        doc_batch_size: Number of documents to process in parallel
        doc_batch_concurrency: Maximum concurrent document processing threads
    """

    regions_dir: str = "media/regions"  # Base directory for extracted images
    do_ocr: bool = False
    do_table_structure: bool = True  # Enable table structure detection
    images_scale: float = 2.0  # Scale factor for generated images
    
    # Maximum number of pages the RT-DETR (layout) model processes in parallel in a single inference pass
    layout_batch_size: int = 1
    # Maximum number of table images (cropped images) that the TableFormer (table-structure) model
    table_batch_size: int = 1

    # setting for batch processing of Docling document restructuring
    # Number of documents that DocumentConverter loads into memory and processes at once
    doc_batch_size: int = 1
    # Number of workers to run concurrently when processing doc_batch_size documents
    # setting this to 1
    doc_batch_concurrency: int = 1

class DoclingParser:
    """
    Document parser that converts various formats to Markdown using Docling library.

    Features:
    - Page-wise processing for paginated documents (PDF)
    - Image extraction and organization by page
    - Support for both paginated (PDF) and linear (DOCX) documents
    """

    def __init__(self, config: Optional[ParserConfig] = None):
        """
        Initialize the Docling parser with configuration.

        Args:
            config: Parser configuration. If None, uses default ParserConfig
        """

        self.config = config or ParserConfig()

        settings.perf.doc_batch_size = self.config.doc_batch_size
        settings.perf.doc_batch_concurrency = self.config.doc_batch_concurrency
        pipeline_options = self._create_pipeline_options(self.config)

        self.converter = DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options)
            }
        )

        self.regions_dir = Path(self.config.regions_dir)
        self.regions_dir.mkdir(parents=True, exist_ok=True)

    @staticmethod
    def _create_pipeline_options(config: ParserConfig) -> PdfPipelineOptions:
        """
        Create PDF processing pipeline options from parser configuration.

        Args:
            config: Parser configuration containing pipeline settings

        Returns:
            Configured PdfPipelineOptions instance with OCR, table detection, and batch settings
        """
        options = PdfPipelineOptions()

        options.do_ocr = config.do_ocr
        options.do_table_structure = config.do_table_structure
        options.generate_picture_images = True
        options.images_scale = config.images_scale

        if hasattr(options, "layout_batch_size"):
            options.layout_batch_size = config.layout_batch_size

        if hasattr(options, "table_batch_size"):
            options.table_batch_size = config.table_batch_size

        return options

    def _input_streams(
        self, 
        file_dict: Dict[str, BytesIO]
    ) -> Tuple[List[DocumentStream], Dict[str, bytes]]:
        """
        file_dict -> DocumentStream
        """
        doc_streams = []
        raw_bytes_map = {}

        for filename, stream in file_dict.items():
            stream.seek(0)
            
            file_bytes = stream.read()
            
            raw_bytes_map[filename] = file_bytes
            
            doc_streams.append(
                DocumentStream(name=filename, stream=BytesIO(file_bytes))
            )
            
        return doc_streams, raw_bytes_map

    def parse(
        self,
        file_dict: Dict[str, BytesIO],
    ) -> Dict[str, str]:
        """
        Parse multiple document files to Markdown format with image extraction in batch.

        Extracts images to organized folders and generates markdown with image references.
        Automatically detects document type and applies appropriate processing.
        Uses batch processing for improved performance.

        Args:
            file_dict: Dictionary mapping filenames (with extensions) to BytesIO file objects

        Returns:
            Dictionary mapping filenames to their converted markdown text content
        """
        # 1. Prepare data (stream list + raw bytes map)
        doc_streams, raw_bytes_map = self._input_streams(file_dict)

        results_map = {}

        # 2. Batch conversion (parallel processing)
        # raises_on_error=False to prevent individual file failures from stopping the entire batch
        conv_results_iter = self.converter.convert_all(doc_streams, raises_on_error=False)

        # 3. Process results
        for result in conv_results_iter:
            filename = result.input.file.name

            if result.status.name != "SUCCESS":
                print(f"[Failed] {filename}: Error during conversion")
                continue

            try:
                # Get original bytes for this file
                file_bytes = raw_bytes_map.get(filename)

                # Create BytesIO object for chart extraction
                file_obj = BytesIO(file_bytes) if file_bytes else None

                # Convert to markdown text
                markdown_text = self._convert_to_markdown(
                    doc=result.document,
                    display_name=filename,
                    file_obj=file_obj
                )

                # Store result
                results_map[filename] = markdown_text
                print(f"[Completed] {filename}")

            except Exception as e:
                print(f"[Error] {filename} post-processing failed: {e}")
                import traceback
                traceback.print_exc()

        return results_map

    def _convert_to_markdown(
        self,
        doc,
        display_name: str,
        file_obj: BytesIO
    ) -> str:
        """
        Convert parsed document to Markdown text based on file type.

        Dispatches to format-specific processing methods:
        - PDF: Hybrid processing (Fitz for images/text + Docling for tables)
        - PPTX: Slide-by-slide with chart extraction (file opened once)
        - Excel: Sheet-by-sheet with chart extraction (file opened once)
        - DOCX: Linear document with image extraction
        - Others: Basic Docling processing

        Args:
            doc: Docling document object containing parsed content
            display_name: File name with extension (e.g., "report.pdf")
            file_obj: BytesIO object containing file data

        Returns:
            Complete markdown text with appropriate headers and content
        """
        path_obj = Path(display_name)
        ext = path_obj.suffix.lower()
        file_key = path_obj.stem

        markdown_text = ""

        # 1. PDF: Hybrid processing (Fitz + Docling)
        if ext == '.pdf':
            markdown_text = self._process_pdf_document(doc, file_key, file_obj)

        # 2. PPTX: Chart extraction + slide-by-slide processing
        elif ext in ['.pptx', '.ppt']:
            markdown_text = self._process_pptx_document(doc, file_key, file_obj)

        # 3. Excel: Sheet names + chart extraction in single pass
        elif ext in ['.xlsx', '.xls', '.xlsm']:
            markdown_text = self._process_excel_document(doc, file_key, file_obj)

        # 4. Word: Linear document with image extraction
        elif ext in ['.docx', '.doc']:
            markdown_text = self._process_docx_document(doc, file_key)

        # 5. Others: Basic processing (Fallback)
        else:
            print(f"[Info] Unknown format {ext}, using default processing.")
            markdown_text = self._process_docx_document(doc, file_key)

        return markdown_text.strip()

    def _process_paginated_document(
        self,
        doc,
        page_num: int,
        file_key: str,
    ) -> str:
        """
        Process a single page from a paginated document (e.g., PDF, Excel).

        Saves page images to page-specific folder (page_XXXX/) and generates markdown for that page.

        Args:
            doc: Docling document object
            page_num: Page number to process (1-indexed)
            file_key: Unique identifier for organizing images

        Returns:
            Markdown text for the specified page with image references
        """
        page_folder = self.regions_dir / file_key / f"page_{page_num:04d}"
        page_folder.mkdir(parents=True, exist_ok=True)

        # Save images for this page and update references
        print(f"  [Saving] Images to {page_folder}")
        page_doc = doc._with_pictures_refs(
            image_dir=page_folder,
            page_no=page_num,
            reference_path=None,
        )

        # Export page to markdown
        markdown_output = page_doc.export_to_markdown(
            page_no=page_num, image_mode=ImageRefMode.REFERENCED
        )

        # Check what files were actually saved
        saved_files = list(page_folder.glob("*")) if page_folder.exists() else []
        print(f"  [Saved] {len(saved_files)} files: {[f.name for f in saved_files]}")

        return markdown_output

    def _process_pdf_document(self, doc, file_key: str, file_obj: BytesIO) -> str:
        """PDF"""
        markdown_parts = []
        num_pages = doc.num_pages()
        
        for page_num in range(1, num_pages + 1):
            page_text = self._process_paginated_document(doc, page_num, file_key)
            
            page_header = f"\n\n- Page {page_num} -\n\n"
            markdown_parts.append(page_header)
            markdown_parts.append(page_text.strip())
            
        return "".join(markdown_parts)

    def _process_docx_document(self, doc, file_key: str) -> str:
        """
        Process linear (non-paginated) documents like DOCX, TXT, etc.

        Saves all images to a single images/ folder and generates markdown.
        Used as fallback for unknown document formats.

        Args:
            doc: Docling document object
            file_key: Unique identifier for organizing images

        Returns:
            Markdown text with image references
        """
        images_dir = self.regions_dir / file_key / "images"
        images_dir.mkdir(parents=True, exist_ok=True) #!

        # Save images and update references
        print(f"  [Saving] Images to {images_dir}") #!
        new_doc = doc._with_pictures_refs(
            image_dir=images_dir, page_no=None, reference_path=None
        )

        markdown_output = new_doc.export_to_markdown(image_mode=ImageRefMode.REFERENCED)

        # Check what files were actually saved
        saved_files = list(images_dir.glob("*")) if images_dir.exists() else [] #!
        print(f"  [Saved] {len(saved_files)} files: {[f.name for f in saved_files]}") #!

        return markdown_output


    def _process_pptx_document(self, doc, file_key: str, file_obj: BytesIO) -> str:
        """
        Process PowerPoint documents with chart extraction.

        Opens file ONCE to extract all chart data, then processes each slide
        with Docling-generated text and inserts charts at appropriate positions.

        Args:
            doc: Docling document object
            file_key: Unique identifier for organizing images
            file_obj: BytesIO object containing PPTX file

        Returns:
            Complete markdown text with slide headers and chart data
        """
        # Open file ONCE to extract all charts
        charts_by_page = self._extract_charts_from_pptx(file_obj)

        markdown_parts = []
        num_pages = doc.num_pages()

        for page_num in range(1, num_pages + 1):
            # Generate base text from Docling
            page_text = self._process_paginated_document(doc, page_num, file_key)

            # Insert charts if present on this slide
            if page_num in charts_by_page:
                page_text = self._insert_charts_at_position(page_text, charts_by_page[page_num])

            page_header = f"\n\n- Slide {page_num} -\n\n"
            markdown_parts.append(page_header)
            markdown_parts.append(page_text.strip())

        return "".join(markdown_parts)

    def _process_excel_document(self, doc, file_key: str, file_obj: BytesIO) -> str:
        """
        Process Excel documents with sheet names and chart extraction.

        Opens file ONCE to extract both sheet names and chart data,
        then processes each sheet with Docling-generated text.

        Args:
            doc: Docling document object
            file_key: Unique identifier for organizing images
            file_obj: BytesIO object containing Excel file

        Returns:
            Complete markdown text with sheet headers and chart data
        """
        # Open file ONCE to extract sheet names AND charts
        sheet_names, charts_by_page = self._extract_excel_metadata(file_obj)

        markdown_parts = []
        num_pages = doc.num_pages()

        for page_num in range(1, num_pages + 1):
            # Generate base text from Docling
            page_text = self._process_paginated_document(doc, page_num, file_key)

            # Create header with sheet name if available
            if (page_num - 1) < len(sheet_names):
                sheet_name = sheet_names[page_num - 1]
                header = f"\n\n- Sheet: {sheet_name} -\n\n"
            else:
                header = f"\n\n- Sheet {page_num} -\n\n"

            # Insert charts if present on this sheet
            if page_num in charts_by_page:
                page_text = self._insert_charts_at_position(page_text, charts_by_page[page_num])

            markdown_parts.append(header)
            markdown_parts.append(page_text.strip())

        return "".join(markdown_parts)

    # ==================== Chart Insertion and Extraction Methods ====================

    def _find_chart_insert_position(self, text: str, pre_text: str, post_text: str, start_pos: int) -> int:
        """
        Find optimal position to insert a chart based on surrounding text.
        """
        if not pre_text and not post_text:
            return -1

        search_text = text[start_pos:]
        
        def normalize(s):  # all line breaks into single spaces for comparison
            return re.sub(r'\s+', ' ', s).strip()

        norm_pre = normalize(pre_text)
        norm_post = normalize(post_text)
        norm_search = normalize(search_text)

        # 1. pre_text exists
        if norm_pre and norm_pre in norm_search:
            idx = text.find(pre_text, start_pos)
            if idx != -1:
                return idx + len(pre_text)
            
            short_pre = pre_text[:20]
            idx = text.find(short_pre, start_pos)
            if idx != -1:
                return idx + len(short_pre)

        # 2. post_text exists (Insert Before)
        if norm_post:
             idx = text.find(post_text, start_pos)
             if idx != -1:
                 return idx
             
             short_post = post_text[:20]
             idx = text.find(short_post, start_pos)
             if idx != -1:
                 return idx

        return -1

    def _insert_charts_at_position(self, markdown_text: str, charts: List[Dict[str, str]]) -> str:
        """
        Insert charts at their original positions in the markdown text.

        Uses context text (pre_text and post_text) to locate the best insertion point.
        Falls back to appending at the end if no matching position is found.

        Args:
            markdown_text: Original markdown text from Docling
            charts: List of chart info dicts with pre_text, post_text, title, and table

        Returns:
            Modified markdown with charts inserted at appropriate positions
        """
        result = markdown_text
        last_pos = 0

        for chart in charts:
            pre_text = chart["pre_text"].strip()
            post_text = chart["post_text"].strip()
            chart_markdown = f"\n# 차트: {chart['title']}\n{chart['table']}\n"

            # Find optimal insertion position
            insert_idx = self._find_chart_insert_position(result, pre_text, post_text, last_pos)

            if insert_idx != -1:
                # Insert chart at found position
                result = result[:insert_idx] + chart_markdown + result[insert_idx:]
                last_pos = insert_idx + len(chart_markdown)
                print(f"# Chart inserted: {chart['title']} (position: {insert_idx})")
            else:
                # Fallback: append to end
                result += chart_markdown
                print(f"# Chart appended: {chart['title']} (no matching position found)")

        return result

    def _get_chart_context(self, shape, all_shapes) -> Tuple[str, str]:
        """
        Find the closest text above and below a chart.

        Args:
            shape: Chart shape object
            all_shapes: List of all shapes in the slide

        Returns:
            Tuple of (text_above, text_below)
        """
        chart_top = shape.top
        chart_bottom = shape.top + shape.height

        text_above = []
        text_below = []

        for other in all_shapes:
            # Skip if it's the chart itself or has no text
            if other == shape or not other.has_text_frame:
                continue

            text_content = other.text_frame.text.strip()
            if not text_content:
                continue

            other_top = other.top
            other_bottom = other.top + other.height

            # Text above the chart
            if other_bottom < chart_top:
                distance = chart_top - other_bottom
                text_above.append((distance, text_content))

            # Text below the chart
            elif other_top > chart_bottom:
                distance = other_top - chart_bottom
                text_below.append((distance, text_content))

        # Sort by distance (closest first)
        text_above.sort(key=lambda x: x[0])
        text_below.sort(key=lambda x: x[0])

        # Return closest text (empty string if none found)
        context_pre = text_above[0][1] if text_above else ""
        context_post = text_below[0][1] if text_below else ""

        return context_pre, context_post

    def _extract_charts_from_pptx(self, file_obj: BytesIO) -> Dict[int, List[Dict[str, str]]]:
        """
        Extract charts from PPTX file and convert to markdown format.

        Args:
            file_obj: BytesIO object containing PPTX file

        Returns:
            Dictionary mapping page numbers to list of chart info dicts with keys:
                - title: Chart title
                - pre_text: Text above the chart
                - post_text: Text below the chart
                - table: Markdown table of chart data
        """
        charts_by_page = {}

        try:
            file_obj.seek(0)
            prs = Presentation(file_obj)

            for slide_idx, slide in enumerate(prs.slides):
                page_num = slide_idx + 1
                all_shapes = list(slide.shapes)

                for shape in slide.shapes:
                    if shape.has_chart:
                        chart = shape.chart

                        # Extract chart title
                        try:
                            title = chart.chart_title.text_frame.text
                        except:
                            title = "no title"

                        # Find context text above and below
                        pre_text, post_text = self._get_chart_context(shape, all_shapes)

                        # Extract chart data and convert to markdown table
                        try:
                            df = pd.DataFrame()
                            plot = chart.plots[0]
                            cats = [c.label for c in plot.categories]

                            # Generate categories if none exist
                            if not cats:
                                cats = [f"Item {i}" for i in range(len(plot.series[0].values))]

                            df.index = cats
                            for ser in plot.series:
                                df[ser.name] = pd.Series(ser.values, index=cats)
                            md_table = df.to_markdown()
                        except Exception as e:
                            md_table = f"(fail to extract data: {str(e)})"

                        # Store chart info
                        chart_info = {
                            "title": title,
                            "pre_text": pre_text,
                            "post_text": post_text,
                            "table": md_table
                        }

                        # Add to page's chart list
                        if page_num not in charts_by_page:
                            charts_by_page[page_num] = []
                        charts_by_page[page_num].append(chart_info)

            return charts_by_page

        except Exception as e:
            print(f"Error extracting charts from PPTX: {str(e)}")
            return {}

    # ==================== Excel Metadata Extraction Methods ====================

    def _extract_excel_metadata(self, file_obj: BytesIO) -> Tuple[List[str], Dict[int, List[Dict[str, str]]]]:
        """
        Extract sheet names AND charts from Excel file in a single pass.

        Opens the Excel file ONCE and extracts all necessary metadata:
        - Sheet names for headers
        - Chart data with position context

        Args:
            file_obj: BytesIO object containing Excel file

        Returns:
            Tuple of (sheet_names, charts_by_page) where:
                - sheet_names: List of sheet names in order
                - charts_by_page: Dictionary mapping page numbers (sheet numbers) to list of chart info dicts
        """
        sheet_names = []
        charts_by_page = {}

        try:
            file_obj.seek(0)
            # data_only=True: Get calculated values instead of formulas
            wb = load_workbook(file_obj, data_only=True)
            sheet_names = wb.sheetnames

            for page_idx, sheet_name in enumerate(sheet_names):
                sheet = wb[sheet_name]
                page_num = page_idx + 1  # 1-based page number

                # Check if sheet has charts
                charts = getattr(sheet, "charts", []) or getattr(sheet, "_charts", [])

                if not charts:
                    continue

                page_charts = []

                for chart in charts:
                    # 1. Extract chart title
                    title = "Untitled Chart"
                    try:
                        if chart.title:
                            # Title entered directly
                            if hasattr(chart.title, 'tx') and chart.title.tx.rich:
                                title = chart.title.tx.rich.p[0].r[0].t
                            # Title referencing a cell
                            elif hasattr(chart.title, 'tx') and chart.title.tx.strRef:
                                ref_vals = self._get_values_from_ref(wb, chart.title.tx.strRef.f)
                                if ref_vals:
                                    title = ref_vals[0]
                    except Exception:
                        pass

                    # 2. Extract context text (Pre/Post Text) based on chart position
                    pre_text, post_text = self._get_chart_context_excel(sheet, chart)

                    # 3. Extract chart data and convert to markdown
                    md_table = self._resolve_excel_chart_data(wb, chart)

                    page_charts.append({
                        "title": title,
                        "pre_text": pre_text,
                        "post_text": post_text,
                        "table": md_table
                    })

                if page_charts:
                    charts_by_page[page_num] = page_charts

            return sheet_names, charts_by_page

        except Exception as e:
            print(f"Error extracting Excel metadata: {str(e)}")
            return [], {}

    def _get_chart_context_excel(self, sheet, chart) -> Tuple[str, str]:
        """
        Extract text above and below a chart in Excel.

        Args:
            sheet: Excel worksheet object
            chart: Chart object

        Returns:
            Tuple of (text_above, text_below)
        """
        pre_text = ""
        post_text = ""
        try:
            # Get anchor information (TwoCellAnchor recommended)
            anchor = chart.anchor
            if hasattr(anchor, '_from'):
                row_start = anchor._from.row
                col_start = anchor._from.col

                # Pre-text: Check row above chart start
                if row_start > 0:
                    cell = sheet.cell(row=row_start, column=col_start + 1)
                    pre_text = str(cell.value) if cell.value else ""

                # Post-text: Check row below chart end
                row_end = anchor.to.row
                cell = sheet.cell(row=row_end + 2, column=col_start + 1)
                post_text = str(cell.value) if cell.value else ""
        except Exception:
            pass  # Skip if position info unavailable or OneCellAnchor

        return pre_text.strip(), post_text.strip()

    def _resolve_excel_chart_data(self, wb, chart) -> str:
        """
        Parse cell references from chart object and convert actual data
        from Pandas DataFrame to Markdown format.

        Args:
            wb: Workbook object
            chart: Chart object

        Returns:
            Markdown table string
        """
        try:
            data_dict = {}
            categories = []

            # --- 1. Get X-axis (category) data ---
            # Usually use first series category reference as common X-axis
            if len(chart.series) > 0:
                cat_ref = None
                try:
                    # Check string reference (strRef) or numeric reference (numRef)
                    first_series = chart.series[0]
                    if hasattr(first_series, 'cat') and first_series.cat:
                        if hasattr(first_series.cat, 'strRef') and first_series.cat.strRef:
                            cat_ref = first_series.cat.strRef.f
                        elif hasattr(first_series.cat, 'numRef') and first_series.cat.numRef:
                            cat_ref = first_series.cat.numRef.f
                    
                    if cat_ref:
                        categories = self._get_values_from_ref(wb, cat_ref)
                except Exception:
                    pass  # Auto-generate if category extraction fails

            # --- 2. Get Y-axis (series values) data ---
            for idx, series in enumerate(chart.series):
                # Series name extract
                series_name = f"Series {idx+1}" 
                try:
                    if series.title:
                        if hasattr(series.title, 'tx') and hasattr(series.title.tx, 'rich'):
                            p_list = getattr(series.title.tx.rich, 'p', [])
                            if p_list:
                                r_list = getattr(p_list[0], 'r', [])
                                if r_list:
                                    series_name = r_list[0].t
                        
                        elif hasattr(series.title, 'tx') and hasattr(series.title.tx, 'strRef'):
                            ref_vals = self._get_values_from_ref(wb, series.title.tx.strRef.f)
                            if ref_vals:
                                series_name = str(ref_vals[0])
                except Exception:
                    pass 

                # Actual data values
                vals = []
                try:
                    if series.val:
                        if hasattr(series.val, 'numRef') and series.val.numRef:
                            vals = self._get_values_from_ref(wb, series.val.numRef.f)
                except Exception:
                    pass

                if vals:
                    data_dict[series_name] = vals

            # --- 3. Convert to Markdown via Pandas ---
            if not data_dict:
                return "(No chart data reference found)"

            # Match data lengths (based on longest data)
            max_len = max(len(v) for v in data_dict.values())

            # Auto-generate categories if missing or length mismatch
            if not categories or len(categories) != max_len:
                categories = [str(i) for i in categories]
                if len(categories) < max_len:
                    categories.extend([f"Item {i+1}" for i in range(len(categories), max_len)])
                else:
                    categories = categories[:max_len]

            for k, v in data_dict.items():
                if len(v) < max_len:
                    data_dict[k] = list(v) + [""] * (max_len - len(v))

            # Create DataFrame
            df = pd.DataFrame(data_dict, index=categories[:max_len])

            return df.to_markdown()

        except Exception as e:
            return f"(Error parsing chart data: {str(e)})"

    def _get_values_from_ref(self, wb, ref_str: str) -> List[Any]:
        """
        Parse reference string like 'Sheet1!$A$1:$A$5' and return actual values.

        Args:
            wb: Workbook object
            ref_str: Cell reference string

        Returns:
            List of cell values
        """
        try:
            if "!" not in ref_str:
                return []

            sheet_part, cell_part = ref_str.rsplit("!", 1)
            # Remove quotes from sheet name ('Sheet 1' -> Sheet 1)
            sheet_name = sheet_part
            if sheet_name.startswith("'") and sheet_name.endswith("'"):
                sheet_name = sheet_name[1:-1]  # 양쪽 따옴표 제거
                sheet_name = sheet_name.replace("''", "'")

            if sheet_name not in wb.sheetnames:
                found = False
                for s in wb.sheetnames:
                    if s.lower() == sheet_name.lower():
                        sheet_name = s
                        found = True
                        break
                if not found:
                    return []

            sheet = wb[sheet_name]

            # Parse range 
            try:
                min_col, min_row, max_col, max_row = range_boundaries(cell_part)
            except ValueError:
                return []

            values = []

            # Read cell values in range order
            for row in sheet.iter_rows(min_row=min_row, max_row=max_row,
                                       min_col=min_col, max_col=max_col,
                                       values_only=True):
                for cell in row:
                    # Handle None values (empty cells as 0 or empty string)
                    val = cell if cell is not None else ""
                    values.append(val)

            return values

        except Exception as e:
            print(f"Error parsing reference ({ref_str}): {e}")
            return []


    

class DocTool:
    """
    High-level document processing tool with simplified interface.

    Wraps DoclingParser to provide easy batch processing of multiple documents.
    Handles various document formats (PDF, DOCX, etc.) and converts them to Markdown
    with automatic image extraction and organization.
    """

    def __init__(
        self,
        regions_dir: Optional[str] = None,
        do_ocr: bool = False,
        do_table_structure: bool = True,
    ):
        """
        Initialize the document processing tool.

        Args:
            regions_dir: Base directory for storing extracted images (default: "media/regions")
            do_ocr: Whether to perform OCR on images within documents
            do_table_structure: Whether to detect and preserve table structures
        """
        config = ParserConfig(
            regions_dir=regions_dir or "media/regions",
            do_ocr=do_ocr,
            do_table_structure=do_table_structure,
        )

        self._parser = DoclingParser(config=config)

    def run(self, file_dict: Dict[str, BytesIO]) -> Dict[str, str]:
        """
        Process multiple documents to Markdown format in batch.

        Args:
            file_dict: Dictionary mapping filenames (with extensions) to BytesIO file objects

        Returns:
            Dictionary mapping filenames to their converted markdown text content
        """
        # Pass all files at once for batch processing
        return self._parser.parse(file_dict)


# example
if __name__ == "__main__":
    input_folder = Path("/home/shaush/pdfs")
    output_root = Path("/home/shaush/work/parsed-outputs")
    
    # 파일 목록 가져오기
    file_list = [p.resolve() for p in input_folder.iterdir() if p.is_file()]
    print(f"Found {len(file_list)} files.")

    processor = DocTool(regions_dir=str(output_root))

    file_dict = {}

    for file_path in file_list:
        with open(file_path, "rb") as f:
                file_dict[file_path.name] = BytesIO(f.read())
    start_time = time.perf_counter()
    results = processor.run(file_dict)
    print("총 파싱시간: ",time.perf_counter() - start_time)