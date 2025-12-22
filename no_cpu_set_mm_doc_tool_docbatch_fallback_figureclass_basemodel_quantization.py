# uv add psutil torch pikepdf
"""
Document Parsing Tool with Optimized Multi-GPU/CPU Processing

Performance Optimizations Applied:
1. CPU Thread Control: Environment variables prevent CPU thread contention
2. Dynamic Worker Allocation: calculate_optimal_workers() balances GPU/CPU resources
3. Generator-based Chunking: _generate_chunks_lazy() eliminates pre-processing bottleneck
4. Long-lived Workers: Increased worker_restart_interval (100) reduces model loading overhead
5. Multi-worker per GPU: workers_per_gpu=2 maximizes GPU utilization
6. No CPU Affinity: OS handles CPU scheduling for better portability

Architecture:
- Main process generates chunks on-demand via generator
- Worker processes start immediately and grab chunks as they're ready
- No blocking wait for all chunks to be pre-processed
- GPU workers can process multiple chunks concurrently per GPU
"""

# Environment variables for CPU thread control (prevents CPU thread contention)
import os
os.environ["OMP_NUM_THREADS"] = "4"
os.environ["MKL_NUM_THREADS"] = "4"
os.environ["OPENBLAS_NUM_THREADS"] = "4"
os.environ["NUMEXPR_NUM_THREADS"] = "4"

import shutil
import time
import math
import tempfile
import multiprocessing
import traceback
import gc
import torch
from pathlib import Path
from io import BytesIO
from typing import List, Dict, Any, Tuple, Optional, Generator
from dataclasses import dataclass
from queue import Empty

# PDF Chunking
import pikepdf

# Docling & Models
from pydantic import BaseModel, Field
from docling.document_converter import DocumentConverter, PdfFormatOption
from docling.datamodel.base_models import InputFormat
from docling.datamodel.pipeline_options import PdfPipelineOptions, AcceleratorOptions
from docling.datamodel.settings import settings
from docling_core.types.doc.base import ImageRefMode
from docling_core.types.io import DocumentStream
from docling_core.types.doc import PictureItem
from docling.backend.pypdfium2_backend import PyPdfiumDocumentBackend

from openpyxl import load_workbook
from openpyxl.utils import range_boundaries
from pptx import Presentation
import pandas as pd
import re



class Figure(BaseModel):
    id: str                         # Image ID for placeholder and figure reference
    mime_type: str                  # MIME type (e.g., "image/png", "image/jpg")
    data: str                       # Base64-encoded image data
 
class Document(BaseModel):
    id: str
    text: str
    images: Optional[List[Figure]] = Field(default=None)

@dataclass
class ChunkResult:
    """For internal processing - Chunked results"""
    original_file_id: str
    chunk_index: int
    text: str
    images: List[Figure]
    success: bool
    error_msg: Optional[str] = None


def calculate_optimal_workers(
    num_gpus: int,
    total_cpus: int,
    workers_per_gpu: int = 2,
    reserved_cpus: int = 2
) -> Tuple[int, Optional[List[int]]]:
    """
    Calculate optimal number of worker processes for GPU/CPU hybrid environments.

    Strategy:
    - GPU mode: Spawn multiple workers per GPU to maximize GPU utilization
    - CPU mode: Use available CPUs minus reserved cores for system operations
    - Hybrid: Balance between GPU capability and CPU availability

    Args:
        num_gpus: Number of available GPUs
        total_cpus: Total number of CPU cores available
        workers_per_gpu: Number of worker processes per GPU (default: 2)
        reserved_cpus: Number of CPUs to reserve for system/main process (default: 2)

    Returns:
        Tuple of (num_workers, gpu_ids) where:
        - num_workers: Optimal number of worker processes
        - gpu_ids: List of GPU IDs to assign (None for CPU-only mode)

    Examples:
        >>> calculate_optimal_workers(num_gpus=2, total_cpus=16, workers_per_gpu=2)
        (4, [0, 1, 0, 1])  # 4 workers cycling through 2 GPUs

        >>> calculate_optimal_workers(num_gpus=0, total_cpus=8, reserved_cpus=2)
        (6, None)  # 6 CPU workers (8 - 2 reserved)
    """
    if num_gpus > 0:
        # GPU mode: workers_per_gpu workers for each GPU
        gpu_based_workers = num_gpus * workers_per_gpu

        # CPU constraint: don't spawn more workers than available CPUs
        available_cpus = max(1, total_cpus - reserved_cpus)
        cpu_limited_workers = min(gpu_based_workers, available_cpus)

        num_workers = cpu_limited_workers

        # Assign GPUs in round-robin fashion
        # Example: 2 GPUs, 4 workers -> [0, 1, 0, 1]
        gpu_ids = [i % num_gpus for i in range(num_workers)]

        return num_workers, gpu_ids
    else:
        # CPU-only mode
        available_cpus = max(1, total_cpus - reserved_cpus)
        num_workers = available_cpus

        return num_workers, None


@dataclass
class ParserConfig:
    """
    Configuration dataclass for DoclingParser.

    Attributes:
        do_ocr: Whether to perform OCR on images within documents
        do_table_structure: Whether to detect and preserve table structures
        do_formula_enrichment: Whether to enable formula enrichment
        generate_picture_images: Whether to generate picture images from documents
        images_scale: Scale factor for image resolution (higher = better quality)
        layout_batch_size: Batch size for layout detection processing
        table_batch_size: Batch size for table structure recognition
        doc_batch_size: Number of documents to process in parallel
        doc_batch_concurrency: Maximum concurrent document processing threads
        chunk_page_size: Number of pages per chunk for PDF splitting
        worker_restart_interval: Number of chunks before worker restart (anti-leak)
        cpu_workers: Number of worker processes when no GPU is available
    """

    do_ocr: bool = False
    do_table_structure: bool = True  # Enable table structure detection
    do_formula_enrichment: bool = True  # Enable formula enrichment
    generate_picture_images: bool = True  # Enable picture image generation
    images_scale: float = 2.0  # Scale factor for generated images

    # Model Batch Sizes (GPU Inference Batch)
    layout_batch_size: int = 16  # Layout detection model batch size
    table_batch_size: int = 16   # Table structure recognition model batch size

    # Document processing settings
    doc_batch_size: int = 8           # Number of documents/chunks processed in batch
    doc_batch_concurrency: int = 1    # Number of concurrent workers in a process (set to 1 for stability)

    # Batch & Resource Settings
    chunk_page_size: int = 10           # Number of pages per chunk
    worker_restart_interval: int = 100  # Restart worker after processing N chunks (Anti-Leak, increased for long-lived workers)

    # Worker allocation settings
    workers_per_gpu: int = 2            # Number of worker processes per GPU (GPU mode)
    reserved_cpus: int = 2              # Number of CPUs reserved for system/main process
    cpu_workers: int = 4                # Fallback: Number of worker processes when no GPU is available (deprecated, use calculate_optimal_workers)


def _split_pdf_to_chunks(
    file_id: str,
    pdf_bytes: bytes,
    chunk_page_size: int
) -> List[Tuple[str, int, BytesIO, int]]:
    """
    [DEPRECATED] Use _generate_chunks_lazy() instead for better performance.

    Split a PDF file into page chunks using pikepdf (preserves ToUnicode maps and image resources).

    This function pre-processes all chunks before returning, creating a bottleneck.
    The generator-based approach (_generate_chunks_lazy) yields chunks on-demand,
    allowing workers to start processing immediately.

    Args:
        file_id: Original filename
        pdf_bytes: PDF file content in bytes
        chunk_page_size: Number of pages per chunk

    Returns:
        List of tuples (chunk_filename, chunk_index, chunk_bytesio, start_page_offset)
    """
    chunks = []

    try:
        # pikepdf can directly open BytesIO
        with pikepdf.open(BytesIO(pdf_bytes)) as pdf:
            total_pages = len(pdf.pages)
            num_chunks = (total_pages + chunk_page_size - 1) // chunk_page_size

            for chunk_idx in range(num_chunks):
                start_page = chunk_idx * chunk_page_size
                end_page = min(start_page + chunk_page_size, total_pages)

                # Create new PDF container
                dst = pikepdf.new()

                # Copy pages (pikepdf preserves resource links during this process)
                for i in range(start_page, end_page):
                    dst.pages.append(pdf.pages[i])

                # Save to memory stream
                chunk_stream = BytesIO()
                dst.save(chunk_stream)
                chunk_stream.seek(0)

                # Generate chunk filename
                chunk_filename = f"{file_id}__chunk_{chunk_idx:04d}.pdf"

                # Include start_page offset for original page numbering
                chunks.append((chunk_filename, chunk_idx, chunk_stream, start_page))

                print(f"[Info] Chunk created: {chunk_filename} (Pages {start_page}-{end_page}, offset: {start_page})")

    except Exception as e:
        print(f"[Error] Failed to split PDF {file_id} with pikepdf: {e}")
        traceback.print_exc()

    return chunks


def _generate_chunks_lazy(
    file_dict: Dict[str, BytesIO],
    chunk_page_size: int
) -> Generator[Tuple[str, int, str, bytes, int, bool, BytesIO], None, None]:
    """
    Generator that yields chunks on-demand for immediate task queue feeding.

    This eliminates the bottleneck where the main process pre-processes all PDFs
    before workers can start. Instead, chunks are generated and fed to workers
    immediately as they become available.

    Args:
        file_dict: Dictionary mapping filenames to BytesIO file objects
        chunk_page_size: Number of pages per PDF chunk

    Yields:
        Tuple of (chunk_filename, chunk_index, original_file_id, chunk_bytes,
                  page_offset, is_pdf, original_stream)
        where:
        - chunk_filename: Name for this chunk
        - chunk_index: Index of this chunk (0 for non-PDF files)
        - original_file_id: Original filename
        - chunk_bytes: Chunk data as bytes
        - page_offset: Starting page number for this chunk
        - is_pdf: Whether this is a PDF file
        - original_stream: Original file stream (for non-PDF files)
    """
    for filename, file_stream in file_dict.items():
        ext = Path(filename).suffix.lower()

        if ext == '.pdf':
            # PDF: Generate chunks on-the-fly
            file_stream.seek(0)
            pdf_bytes = file_stream.read()

            try:
                with pikepdf.open(BytesIO(pdf_bytes)) as pdf:
                    total_pages = len(pdf.pages)
                    num_chunks = (total_pages + chunk_page_size - 1) // chunk_page_size

                    print(f"[ChunkGenerator] Streaming {num_chunks} chunks from {filename}")

                    for chunk_idx in range(num_chunks):
                        start_page = chunk_idx * chunk_page_size
                        end_page = min(start_page + chunk_page_size, total_pages)

                        # Create chunk
                        dst = pikepdf.new()
                        for i in range(start_page, end_page):
                            dst.pages.append(pdf.pages[i])

                        # Save to memory
                        chunk_stream = BytesIO()
                        dst.save(chunk_stream)
                        chunk_stream.seek(0)
                        chunk_bytes = chunk_stream.read()

                        chunk_filename = f"{filename}__chunk_{chunk_idx:04d}.pdf"

                        # Yield immediately - worker can grab this right away!
                        yield (
                            chunk_filename,
                            chunk_idx,
                            filename,  # original_file_id
                            chunk_bytes,
                            start_page,  # page_offset
                            True,  # is_pdf
                            None  # original_stream (not needed for chunked PDFs)
                        )

            except Exception as e:
                print(f"[ChunkGenerator] Error processing {filename}: {e}")
                traceback.print_exc()

        else:
            # Non-PDF: Yield as-is
            file_stream.seek(0)
            file_bytes = file_stream.read()

            yield (
                filename,
                0,  # chunk_index
                filename,  # original_file_id
                file_bytes,
                0,  # page_offset
                False,  # is_pdf
                file_stream  # original_stream
            )


class DoclingParser:
    """
    Document parser that converts various formats to Markdown using Docling library.

    Features:
    - Page-wise processing for paginated documents (PDF)
    - Image extraction and organization by page
    - Support for both paginated (PDF) and linear (DOCX) documents
    """

    def __init__(self, config: Optional[ParserConfig] = None, gpu_id: Optional[int] = None):
        """
        Initialize the Docling parser with configuration.

        Args:
            config: Parser configuration. If None, uses default ParserConfig
            gpu_id: GPU device ID to use for processing (None for CPU or default GPU)
        """

        self.config = config or ParserConfig()
        self.gpu_id = gpu_id

        settings.perf.doc_batch_concurrency = self.config.doc_batch_concurrency
        pipeline_options = self._create_pipeline_options(self.config)

        # Configure GPU if specified
        if gpu_id is not None:
            pipeline_options.accelerator_options = AcceleratorOptions(
                num_threads=4,
                device=f"cuda:{gpu_id}"
            )

        # Primary converter with default backend
        # Current PDF Backend: <class 'docling.backend.docling_parse_v4_backend.DoclingParseV4DocumentBackend'>
        self.converter = DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(
                    pipeline_options=pipeline_options
                )
            }
        )
        # PDF 포맷을 담당하는 백엔드 객체 가져오기
        pdf_backend = self.converter.format_to_options.get(InputFormat.PDF)
        print(f"Current PDF Backend: {pdf_backend.backend}")


        # Fallback converter with PyPdfiumDocumentBackend for handling "Invalid code point" errors
        self.fallback_converter = DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(
                    pipeline_options=pipeline_options,
                    backend=PyPdfiumDocumentBackend
                )
            }
        )

    @staticmethod
    def _create_pipeline_options(config: ParserConfig) -> PdfPipelineOptions:
        """
        Create PDF processing pipeline options from parser configuration.

        Args:
            config: Parser configuration containing pipeline settings

        Returns:
            Configured PdfPipelineOptions instance with OCR, table detection, and batch settings
        """
        options = PdfPipelineOptions()

        options.do_ocr = config.do_ocr
        options.do_table_structure = config.do_table_structure
        options.do_formula_enrichment = config.do_formula_enrichment
        options.generate_picture_images = config.generate_picture_images
        options.images_scale = config.images_scale

        if hasattr(options, "layout_batch_size"):
            options.layout_batch_size = config.layout_batch_size

        if hasattr(options, "table_batch_size"):
            options.table_batch_size = config.table_batch_size

        return options

    def _input_streams(
        self, 
        file_dict: Dict[str, BytesIO]
    ) -> Tuple[List[DocumentStream], Dict[str, bytes]]:
        """file_dict -> DocumentStream"""
        doc_streams = []
        raw_bytes_map = {}

        for filename, stream in file_dict.items():
            stream.seek(0)
            
            file_bytes = stream.read() 
            
            raw_bytes_map[filename] = file_bytes
            
            doc_streams.append(
                DocumentStream(name=filename, stream=BytesIO(file_bytes))
            )
            
        return doc_streams, raw_bytes_map

    def parse(
        self,
        file_dict: Dict[str, BytesIO],
        page_offset: int = 0
    ) -> Dict[str, str]:
        """
        Parse multiple document files to Markdown format with image extraction in batch.

        Extracts images to organized folders and generates markdown with image references.
        Automatically detects document type and applies appropriate processing.
        Uses batch processing for improved performance.

        Implements automatic fallback to PyPdfiumDocumentBackend when "Invalid code point"
        errors occur during parsing.

        Args:
            file_dict: Dictionary mapping filenames (with extensions) to BytesIO file objects
            page_offset: Page number offset for PDF chunks (maintains original page numbers)

        Returns:
            Dictionary mapping filenames to their converted markdown text content
        """
        # Prepare input streams
        doc_streams, raw_bytes_map = self._input_streams(file_dict)
        results_map = {}

        # Execute primary batch conversion
        # raises_on_error = False : the iterator yields failure results instead of crashing.
        primary_iter = self.converter.convert_all(doc_streams, raises_on_error=False)

        for result in primary_iter:
            filename = result.input.file.name

            # Primary conversion successful
            if result.status.name == "SUCCESS":
                self._finalize_result(result, filename, raw_bytes_map, results_map, page_offset)
                continue

            # Primary failed -> Analyze errors
            # Check for "Invalid code point" error.
            is_target_error = False
            for err in result.errors:
                if "Invalid code point" in str(err.error_message):
                    is_target_error = True
                    break

            # Trigger immediate fallback if critical error detected
            if is_target_error:
                
                try:
                    file_bytes = raw_bytes_map.get(filename)
                    if not file_bytes:
                        print(f"[Error] Bytes missing for {filename}")
                        continue

                    retry_stream = DocumentStream(name=filename, stream=BytesIO(file_bytes))

                    # Run fallback converter on the single file
                    fallback_iter = self.fallback_converter.convert_all([retry_stream], raises_on_error=False)
                    retry_result = next(fallback_iter)

                    if retry_result.status.name == "SUCCESS":
                        print(f"[Fallback Success] Recovered {filename}")
                        self._finalize_result(retry_result, filename, raw_bytes_map, results_map, page_offset)
                    else:
                        print(f"  [Fallback Failed] {filename} failed again.")
                        for e in retry_result.errors:
                            print(f"   - Error: {e.error_message}")

                except Exception as e:
                    print(f"  [Fallback Critical] Error during PyPdfiumBackend: {e}")
                    import traceback
                    traceback.print_exc()
            
            # Non-recoverable error
            else:
                print(f"  [Failure] {filename} failed with error:")
                for err in result.errors:
                    print(f"   - {err.error_message}")

        return list(results_map.values())

    def _finalize_result(self, result, filename, raw_bytes_map, results_map, page_offset=0):
        try:
            # Prepare file object for formats that need it
            file_obj = None
            ext = Path(filename).suffix.lower()
            if ext in ['.pptx', '.ppt', '.xlsx', '.xls', '.xlsm']:
                file_bytes = raw_bytes_map.get(filename)
                if file_bytes:
                    file_obj = BytesIO(file_bytes)

            # Extract markdown text and figures
            markdown_text, figures = self._convert_to_document_content(
                doc=result.document,
                display_name=filename,
                file_obj=file_obj,
                page_offset=page_offset
            )

            # Create Document object and store in results map
            doc_obj = Document(
                id=filename,
                text=markdown_text,
                images=figures
            )

            results_map[filename] = doc_obj
            
            print(f"[Completed] {filename} (Extracted {len(figures)} images)")

        except Exception as e:
            print(f"[Error] Finalizing {filename}: {e}")
            import traceback
            traceback.print_exc()

    def _extract_figures_and_patch_doc(self, doc, file_key: str, page_offset: int = 0) -> List[Figure]:
        figures = []
        for item, _ in doc.iterate_items():
            if isinstance(item, PictureItem):
                img = item.get_image(doc=doc)
                if img:
                    # Generate unique image ID with original page numbering
                    page_no = item.prov[0].page_no if item.prov else 0
                    actual_page_no = page_no + page_offset  # Apply offset for chunks
                    self_ref = item.self_ref.replace("#/", "").replace("/", "_")
                    if page_no == 0:
                        img_id = f"{file_key}/images/{self_ref}.png"
                    else:
                        actual_page_no_str = f"{actual_page_no:04d}"
                        img_id = f"{file_key}/page_{actual_page_no_str}/{self_ref}.png"

                    # Create Figure object with Base64 data
                    figures.append(Figure(
                        id=img_id,
                        mime_type="image/png",
                        data=item._image_to_base64(img)
                    ))

                    # Patch document object for markdown links
                    item.image.uri = Path(img_id)
        return figures


    def _convert_to_document_content(
        self,
        doc,
        display_name: str,
        file_obj: Optional[BytesIO] = None,
        page_offset: int = 0
    ) -> Tuple[str, List[Figure]]:

        path_obj = Path(display_name)
        ext = path_obj.suffix.lower()
        file_key = path_obj.stem

        # Handlers return (text, figures) tuple
        handlers = {
            '.pdf': lambda: self._process_pdf_document(doc, file_key, file_obj, page_offset),
            '.pptx': lambda: self._process_pptx_document(doc, file_key, file_obj, page_offset),
            '.ppt': lambda: self._process_pptx_document(doc, file_key, file_obj, page_offset),
            '.xlsx': lambda: self._process_excel_document(doc, file_key, file_obj, page_offset),
            '.xls': lambda: self._process_excel_document(doc, file_key, file_obj, page_offset),
            '.xlsm': lambda: self._process_excel_document(doc, file_key, file_obj, page_offset),
            '.docx': lambda: self._process_docx_document(doc, file_key, page_offset),
            '.doc': lambda: self._process_docx_document(doc, file_key, page_offset),
        }

        handler = handlers.get(ext, lambda: self._process_docx_document(doc, file_key, page_offset))

        text, figures = handler()
        return text.strip(), figures

    def _process_pdf_document(self, doc, file_key: str, file_obj: BytesIO, page_offset: int = 0) -> Tuple[str, List[Figure]]:
        # Extract images and patch document
        all_figures = self._extract_figures_and_patch_doc(doc, file_key, page_offset)

        markdown_parts = []
        # Generate markdown with REFERENCED mode
        for page_num in range(1, doc.num_pages() + 1):
            page_md = doc.export_to_markdown(
                page_no=page_num,
                image_mode=ImageRefMode.REFERENCED
            )
            # Apply page offset to maintain original page numbering
            actual_page_num = page_num + page_offset
            markdown_parts.append(f"\n\n- Page {actual_page_num} -\n\n{page_md.strip()}")

        return "".join(markdown_parts), all_figures

    def _process_docx_document(self, doc, file_key: str, page_offset: int = 0) -> Tuple[str, List[Figure]]:
        # Extract images and patch document
        figures = self._extract_figures_and_patch_doc(doc, file_key, page_offset)

        # Generate markdown
        text = doc.export_to_markdown(image_mode=ImageRefMode.REFERENCED)
        return text, figures

    def _process_pptx_document(self, doc, file_key: str, file_obj: BytesIO, page_offset: int = 0) -> Tuple[str, List[Figure]]:
        # Extract images and patch document
        all_figures = self._extract_figures_and_patch_doc(doc, file_key, page_offset)

        # Extract charts
        charts_by_page = self._extract_charts_from_pptx(file_obj)

        markdown_parts = []
        for page_num in range(1, doc.num_pages() + 1):
            page_text = doc.export_to_markdown(
                page_no=page_num,
                image_mode=ImageRefMode.REFERENCED
            )

            if page_num in charts_by_page:
                page_text = self._insert_charts_at_position(page_text, charts_by_page[page_num])

            # Apply page offset to maintain original slide numbering
            actual_page_num = page_num + page_offset
            markdown_parts.append(f"\n\n- Slide {actual_page_num} -\n\n{page_text.strip()}")

        return "".join(markdown_parts), all_figures

    def _process_excel_document(self, doc, file_key: str, file_obj: BytesIO, page_offset: int = 0) -> Tuple[str, List[Figure]]:
        # Extract images and patch document
        all_figures = self._extract_figures_and_patch_doc(doc, file_key, page_offset)

        sheet_names, charts_by_page = self._extract_excel_metadata(file_obj)
        markdown_parts = []

        for page_num in range(1, doc.num_pages() + 1):
            page_text = doc.export_to_markdown(
                page_no=page_num,
                image_mode=ImageRefMode.REFERENCED
            )

            # Apply page offset to sheet numbering
            actual_page_num = page_num + page_offset

            # Add sheet header
            if (page_num - 1) < len(sheet_names):
                header = f"\n\n- Sheet: {sheet_names[page_num - 1]} (Page {actual_page_num}) -\n\n"
            else:
                header = f"\n\n- Sheet {actual_page_num} -\n\n"

            if page_num in charts_by_page:
                page_text = self._insert_charts_at_position(page_text, charts_by_page[page_num])

            markdown_parts.append(header)
            markdown_parts.append(page_text.strip())

        return "".join(markdown_parts), all_figures

    # ==================== Chart Insertion and Extraction Methods ====================

    def _find_chart_insert_position(self, text: str, pre_text: str, post_text: str, start_pos: int) -> int:
        """
        Find optimal position to insert a chart based on surrounding text.
        """
        if not pre_text and not post_text:
            return -1

        search_text = text[start_pos:]
        
        def normalize(s):  # Normalize whitespace for comparison
            return re.sub(r'\s+', ' ', s).strip()

        norm_pre = normalize(pre_text)
        norm_post = normalize(post_text)
        norm_search = normalize(search_text)

        # Try to find position using pre_text
        if norm_pre and norm_pre in norm_search:
            idx = text.find(pre_text, start_pos)
            if idx != -1:
                return idx + len(pre_text)
            
            short_pre = pre_text[:20]
            idx = text.find(short_pre, start_pos)
            if idx != -1:
                return idx + len(short_pre)

        # Try to find position using post_text
        if norm_post:
             idx = text.find(post_text, start_pos)
             if idx != -1:
                 return idx

             short_post = post_text[:20]
             idx = text.find(short_post, start_pos)
             if idx != -1:
                 return idx

        return -1

    def _insert_charts_at_position(self, markdown_text: str, charts: List[Dict[str, str]]) -> str:
        """
        Insert charts at their original positions in the markdown text.

        Uses context text (pre_text and post_text) to locate the best insertion point.
        Falls back to appending at the end if no matching position is found.

        Args:
            markdown_text: Original markdown text from Docling
            charts: List of chart info dicts with pre_text, post_text, title, and table

        Returns:
            Modified markdown with charts inserted at appropriate positions
        """
        result = markdown_text
        last_pos = 0

        for chart in charts:
            pre_text = chart["pre_text"].strip()
            post_text = chart["post_text"].strip()
            chart_markdown = f"\n# Chart: {chart['title']}\n{chart['table']}\n"

            # Find optimal insertion position
            insert_idx = self._find_chart_insert_position(result, pre_text, post_text, last_pos)

            if insert_idx != -1:
                # Insert chart at found position
                result = result[:insert_idx] + chart_markdown + result[insert_idx:]
                last_pos = insert_idx + len(chart_markdown)
                print(f"# Chart inserted: {chart['title']} (position: {insert_idx})")
            else:
                # Fallback: append to end
                result += chart_markdown
                print(f"# Chart appended: {chart['title']} (no matching position found)")

        return result

    def _get_chart_context(self, shape, all_shapes) -> Tuple[str, str]:
        """
        Find the closest text above and below a chart.

        Args:
            shape: Chart shape object
            all_shapes: List of all shapes in the slide

        Returns:
            Tuple of (text_above, text_below)
        """
        chart_top = shape.top
        chart_bottom = shape.top + shape.height

        text_above = []
        text_below = []

        for other in all_shapes:
            # Skip if it's the chart itself or has no text
            if other == shape or not other.has_text_frame:
                continue

            text_content = other.text_frame.text.strip()
            if not text_content:
                continue

            other_top = other.top
            other_bottom = other.top + other.height

            # Text above the chart
            if other_bottom < chart_top:
                distance = chart_top - other_bottom
                text_above.append((distance, text_content))

            # Text below the chart
            elif other_top > chart_bottom:
                distance = other_top - chart_bottom
                text_below.append((distance, text_content))

        # Sort by distance (closest first)
        text_above.sort(key=lambda x: x[0])
        text_below.sort(key=lambda x: x[0])

        # Return closest text (empty string if none found)
        context_pre = text_above[0][1] if text_above else ""
        context_post = text_below[0][1] if text_below else ""

        return context_pre, context_post

    def _extract_charts_from_pptx(self, file_obj: BytesIO) -> Dict[int, List[Dict[str, str]]]:
        """
        Extract charts from PPTX file and convert to markdown format.

        Args:
            file_obj: BytesIO object containing PPTX file

        Returns:
            Dictionary mapping page numbers to list of chart info dicts with keys:
                - title: Chart title
                - pre_text: Text above the chart
                - post_text: Text below the chart
                - table: Markdown table of chart data
        """
        charts_by_page = {}

        try:
            file_obj.seek(0)
            prs = Presentation(file_obj)

            for slide_idx, slide in enumerate(prs.slides):
                page_num = slide_idx + 1
                all_shapes = list(slide.shapes)

                for shape in slide.shapes:
                    if shape.has_chart:
                        chart = shape.chart

                        # Extract chart title
                        try:
                            title = chart.chart_title.text_frame.text
                        except:
                            title = "no title"

                        # Find context text above and below
                        pre_text, post_text = self._get_chart_context(shape, all_shapes)

                        # Extract chart data and convert to markdown table
                        try:
                            df = pd.DataFrame()
                            plot = chart.plots[0]
                            cats = [c.label for c in plot.categories]

                            # Generate categories if none exist
                            if not cats:
                                cats = [f"Item {i}" for i in range(len(plot.series[0].values))]

                            df.index = cats
                            for ser in plot.series:
                                df[ser.name] = pd.Series(ser.values, index=cats)
                            md_table = df.to_markdown()
                        except Exception as e:
                            md_table = f"(fail to extract data: {str(e)})"

                        # Store chart info
                        chart_info = {
                            "title": title,
                            "pre_text": pre_text,
                            "post_text": post_text,
                            "table": md_table
                        }

                        # Add to page's chart list
                        if page_num not in charts_by_page:
                            charts_by_page[page_num] = []
                        charts_by_page[page_num].append(chart_info)

            return charts_by_page

        except Exception as e:
            print(f"Error extracting charts from PPTX: {str(e)}")
            return {}

    # ==================== Excel Metadata Extraction Methods ====================

    def _extract_excel_metadata(self, file_obj: BytesIO) -> Tuple[List[str], Dict[int, List[Dict[str, str]]]]:
        """
        Extract sheet names AND charts from Excel file in a single pass.

        Opens the Excel file ONCE and extracts all necessary metadata:
        - Sheet names for headers
        - Chart data with position context

        Args:
            file_obj: BytesIO object containing Excel file

        Returns:
            Tuple of (sheet_names, charts_by_page) where:
                - sheet_names: List of sheet names in order
                - charts_by_page: Dictionary mapping page numbers (sheet numbers) to list of chart info dicts
        """
        sheet_names = []
        charts_by_page = {}

        try:
            file_obj.seek(0)
            # data_only=True: Get calculated values instead of formulas
            wb = load_workbook(file_obj, data_only=True)
            sheet_names = wb.sheetnames

            for page_idx, sheet_name in enumerate(sheet_names):
                sheet = wb[sheet_name]
                page_num = page_idx + 1  # 1-based page number

                # Check if sheet has charts
                charts = getattr(sheet, "charts", []) or getattr(sheet, "_charts", [])

                if not charts:
                    continue

                page_charts = []

                for chart in charts:
                    # Extract chart title
                    title = "Untitled Chart"
                    try:
                        if chart.title:
                            # Title entered directly
                            if hasattr(chart.title, 'tx') and chart.title.tx.rich:
                                title = chart.title.tx.rich.p[0].r[0].t
                            # Title referencing a cell
                            elif hasattr(chart.title, 'tx') and chart.title.tx.strRef:
                                ref_vals = self._get_values_from_ref(wb, chart.title.tx.strRef.f)
                                if ref_vals:
                                    title = ref_vals[0]
                    except Exception:
                        pass

                    # Extract context text (Pre/Post Text) based on chart position
                    pre_text, post_text = self._get_chart_context_excel(sheet, chart)

                    # Extract chart data and convert to markdown
                    md_table = self._resolve_excel_chart_data(wb, chart)

                    page_charts.append({
                        "title": title,
                        "pre_text": pre_text,
                        "post_text": post_text,
                        "table": md_table
                    })

                if page_charts:
                    charts_by_page[page_num] = page_charts

            return sheet_names, charts_by_page

        except Exception as e:
            print(f"Error extracting Excel metadata: {str(e)}")
            return [], {}

    def _get_chart_context_excel(self, sheet, chart) -> Tuple[str, str]:
        """
        Extract text above and below a chart in Excel.

        Args:
            sheet: Excel worksheet object
            chart: Chart object

        Returns:
            Tuple of (text_above, text_below)
        """
        pre_text = ""
        post_text = ""
        try:
            # Get anchor information (TwoCellAnchor recommended)
            anchor = chart.anchor
            if hasattr(anchor, '_from'):
                row_start = anchor._from.row
                col_start = anchor._from.col

                # Pre-text: Check row above chart start
                if row_start > 0:
                    cell = sheet.cell(row=row_start, column=col_start + 1)
                    pre_text = str(cell.value) if cell.value else ""

                # Post-text: Check row below chart end
                row_end = anchor.to.row
                cell = sheet.cell(row=row_end + 2, column=col_start + 1)
                post_text = str(cell.value) if cell.value else ""
        except Exception:
            pass  # Skip if position info unavailable or OneCellAnchor

        return pre_text.strip(), post_text.strip()

    def _resolve_excel_chart_data(self, wb, chart) -> str:
        """
        Parse cell references from chart object and convert actual data
        from Pandas DataFrame to Markdown format.

        Args:
            wb: Workbook object
            chart: Chart object

        Returns:
            Markdown table string
        """
        try:
            data_dict = {}
            categories = []

            # Get X-axis (category) data
            # Usually use first series category reference as common X-axis
            if len(chart.series) > 0:
                cat_ref = None
                try:
                    # Check string reference (strRef) or numeric reference (numRef)
                    first_series = chart.series[0]
                    if hasattr(first_series, 'cat') and first_series.cat:
                        if hasattr(first_series.cat, 'strRef') and first_series.cat.strRef:
                            cat_ref = first_series.cat.strRef.f
                        elif hasattr(first_series.cat, 'numRef') and first_series.cat.numRef:
                            cat_ref = first_series.cat.numRef.f
                    
                    if cat_ref:
                        categories = self._get_values_from_ref(wb, cat_ref)
                except Exception:
                    pass  # Auto-generate if category extraction fails

            # Get Y-axis (series values) data
            for idx, series in enumerate(chart.series):
                # Series name extract
                series_name = f"Series {idx+1}" 
                try:
                    if series.title:
                        if hasattr(series.title, 'tx') and hasattr(series.title.tx, 'rich'):
                            p_list = getattr(series.title.tx.rich, 'p', [])
                            if p_list:
                                r_list = getattr(p_list[0], 'r', [])
                                if r_list:
                                    series_name = r_list[0].t
                        
                        elif hasattr(series.title, 'tx') and hasattr(series.title.tx, 'strRef'):
                            ref_vals = self._get_values_from_ref(wb, series.title.tx.strRef.f)
                            if ref_vals:
                                series_name = str(ref_vals[0])
                except Exception:
                    pass 

                # Actual data values
                vals = []
                try:
                    if series.val:
                        if hasattr(series.val, 'numRef') and series.val.numRef:
                            vals = self._get_values_from_ref(wb, series.val.numRef.f)
                except Exception:
                    pass

                if vals:
                    data_dict[series_name] = vals

            # Convert to Markdown via Pandas
            if not data_dict:
                return "(No chart data reference found)"

            # Match data lengths (based on longest data)
            max_len = max(len(v) for v in data_dict.values())

            # Auto-generate categories if missing or length mismatch
            if not categories or len(categories) != max_len:
                categories = [str(i) for i in categories]
                if len(categories) < max_len:
                    categories.extend([f"Item {i+1}" for i in range(len(categories), max_len)])
                else:
                    categories = categories[:max_len]

            for k, v in data_dict.items():
                if len(v) < max_len:
                    data_dict[k] = list(v) + [""] * (max_len - len(v))

            # Create DataFrame
            df = pd.DataFrame(data_dict, index=categories[:max_len])

            return df.to_markdown()

        except Exception as e:
            return f"(Error parsing chart data: {str(e)})"

    def _get_values_from_ref(self, wb, ref_str: str) -> List[Any]:
        """
        Parse reference string like 'Sheet1!$A$1:$A$5' and return actual values.

        Args:
            wb: Workbook object
            ref_str: Cell reference string

        Returns:
            List of cell values
        """
        try:
            if "!" not in ref_str:
                return []

            sheet_part, cell_part = ref_str.rsplit("!", 1)
            # Remove quotes from sheet name ('Sheet 1' -> Sheet 1)
            sheet_name = sheet_part
            if sheet_name.startswith("'") and sheet_name.endswith("'"):
                sheet_name = sheet_name[1:-1]
                sheet_name = sheet_name.replace("''", "'")

            if sheet_name not in wb.sheetnames:
                found = False
                for s in wb.sheetnames:
                    if s.lower() == sheet_name.lower():
                        sheet_name = s
                        found = True
                        break
                if not found:
                    return []

            sheet = wb[sheet_name]

            # Parse range 
            try:
                min_col, min_row, max_col, max_row = range_boundaries(cell_part)
            except ValueError:
                return []

            values = []

            # Read cell values in range order
            for row in sheet.iter_rows(min_row=min_row, max_row=max_row,
                                       min_col=min_col, max_col=max_col,
                                       values_only=True):
                for cell in row:
                    # Handle None values (empty cells as 0 or empty string)
                    val = cell if cell is not None else ""
                    values.append(val)

            return values

        except Exception as e:
            print(f"Error parsing reference ({ref_str}): {e}")
            return []


def _chunk_worker_process(
    worker_id: int,
    gpu_id: Optional[int],
    task_queue: multiprocessing.Queue,
    result_queue: multiprocessing.Queue,
    config_dict: Dict[str, Any],
    worker_restart_interval: int
):
    """
    Worker process that processes document chunks from a queue.

    Args:
        worker_id: Unique worker identifier
        gpu_id: GPU device ID (None for CPU)
        task_queue: Queue containing (chunk_filename, chunk_index, original_file_id, chunk_bytes, page_offset)
        result_queue: Queue for returning ChunkResult objects
        config_dict: Configuration dictionary
        worker_restart_interval: Number of chunks to process before self-termination
    """
    device_str = f"GPU-{gpu_id}" if gpu_id is not None else f"CPU-Worker{worker_id}"
    print(f"  [Worker-{device_str}] Worker process started")

    try:
        # Reconstruct config
        config = ParserConfig(**config_dict)

        # Initialize parser with GPU assignment
        parser = DoclingParser(config=config, gpu_id=gpu_id)

        chunks_processed = 0

        while True:
            try:
                # Get task from queue with timeout
                task = task_queue.get(timeout=5)

                if task is None:  # Poison pill to terminate worker
                    print(f"  [Worker-{device_str}] Received shutdown signal")
                    break

                chunk_filename, chunk_index, original_file_id, chunk_bytes, page_offset = task

                # Start timing
                start_time = time.perf_counter()

                # Process chunk with page offset for correct page numbering
                file_dict = {chunk_filename: BytesIO(chunk_bytes)}
                doc_list = parser.parse(file_dict, page_offset=page_offset)

                # Calculate processing time
                total_time = time.perf_counter() - start_time

                if doc_list and len(doc_list) > 0:
                    doc = doc_list[0]

                    # Create ChunkResult
                    chunk_result = ChunkResult(
                        original_file_id=original_file_id,
                        chunk_index=chunk_index,
                        text=doc.text,
                        images=doc.images if doc.images else [],
                        success=True,
                        error_msg=None
                    )
                else:
                    chunk_result = ChunkResult(
                        original_file_id=original_file_id,
                        chunk_index=chunk_index,
                        text="",
                        images=[],
                        success=False,
                        error_msg="No document returned from parser"
                    )
                    print(f"  [Worker-{device_str}] Chunk {chunk_index} of {original_file_id} failed ({total_time:.2f}s)")

                result_queue.put(chunk_result)
                chunks_processed += 1

                # Self-restart mechanism
                if chunks_processed >= worker_restart_interval:
                    print(f"  [Worker-{device_str}] Shutting down after {chunks_processed} chunks (restart interval reached)")
                    break

            except Empty:
                continue

            except Exception as e:
                print(f"[Worker-{device_str}] Error processing chunk: {e}")
                traceback.print_exc()

                # Return error result
                if 'chunk_index' in locals() and 'original_file_id' in locals():
                    error_result = ChunkResult(
                        original_file_id=original_file_id,
                        chunk_index=chunk_index,
                        text="",
                        images=[],
                        success=False,
                        error_msg=str(e)
                    )
                    result_queue.put(error_result)

    except Exception as e:
        print(f"[Worker-{device_str}] Fatal error in worker: {e}")
        traceback.print_exc()

    finally:
        print(f"  [Worker-{device_str}] Worker process terminated (total processed: {chunks_processed} chunks)")


class WorkerManager:
    """
    Manages worker processes and handles automatic restart on interval.
    """

    def __init__(
        self,
        num_workers: int,
        gpu_ids: Optional[List[int]],
        config_dict: Dict[str, Any],
        worker_restart_interval: int
    ):
        """
        Args:
            num_workers: Number of worker processes
            gpu_ids: List of GPU IDs (None for CPU-only mode)
            config_dict: Configuration dictionary
            worker_restart_interval: Chunks per worker before restart
        """
        self.num_workers = num_workers
        self.gpu_ids = gpu_ids
        self.config_dict = config_dict
        self.worker_restart_interval = worker_restart_interval

        self.task_queue = multiprocessing.Queue()
        self.result_queue = multiprocessing.Queue()

        self.processes: List[multiprocessing.Process] = []

    def start_workers(self):
        """Start all worker processes."""
        for i in range(self.num_workers):
            gpu_id = self.gpu_ids[i] if self.gpu_ids else None
            self._start_single_worker(i, gpu_id)

    def _start_single_worker(self, worker_id: int, gpu_id: Optional[int]):
        """Start a single worker process."""
        p = multiprocessing.Process(
            target=_chunk_worker_process,
            args=(
                worker_id,
                gpu_id,
                self.task_queue,
                self.result_queue,
                self.config_dict,
                self.worker_restart_interval
            )
        )
        p.start()
        self.processes.append(p)

    def restart_worker(self, worker_id: int, gpu_id: Optional[int]):
        """Restart a specific worker."""
        if worker_id < len(self.processes):
            old_process = self.processes[worker_id]
            if old_process.is_alive():
                old_process.terminate()
                old_process.join(timeout=5)

            # Start new worker process
            p = multiprocessing.Process(
                target=_chunk_worker_process,
                args=(
                    worker_id,
                    gpu_id,
                    self.task_queue,
                    self.result_queue,
                    self.config_dict,
                    self.worker_restart_interval
                )
            )
            p.start()
            self.processes[worker_id] = p  # Update process list
            device_name = f"GPU-{gpu_id}" if gpu_id is not None else f"CPU-Worker{worker_id}"
            print(f" [Manager] Worker {device_name} restarted successfully")

    def shutdown(self):
        """Shutdown all workers gracefully."""
        print(f"[Manager] Sending shutdown signal to {self.num_workers} workers...")

        # Send poison pills
        for _ in range(self.num_workers):
            self.task_queue.put(None)

        # Wait for processes to finish
        for p in self.processes:
            p.join(timeout=10)
            if p.is_alive():
                p.terminate()

        print("[Manager] All workers shut down successfully")


class DocTool:
    """
    High-level document processing tool with Multi-GPU/CPU support.

    Wraps DoclingParser to provide easy batch processing with automatic
    chunking and parallel processing across GPUs or CPUs.
    """

    def __init__(
        self,
        do_ocr: bool = False,
        do_table_structure: bool = True,
        chunk_page_size: int = 10,
        worker_restart_interval: int = 100,
        workers_per_gpu: int = 2,
        reserved_cpus: int = 2,
        cpu_workers: int = 4,
    ):
        """
        Initialize the document processing tool.

        Args:
            do_ocr: Whether to perform OCR on images within documents
            do_table_structure: Whether to detect and preserve table structures
            chunk_page_size: Number of pages per PDF chunk
            worker_restart_interval: Number of chunks before worker restart (increased default for long-lived workers)
            workers_per_gpu: Number of worker processes per GPU (default: 2)
            reserved_cpus: Number of CPUs reserved for system/main process (default: 2)
            cpu_workers: Fallback number of CPU worker processes (deprecated, use calculate_optimal_workers)
        """
        self.config = ParserConfig(
            do_ocr=do_ocr,
            do_table_structure=do_table_structure,
            chunk_page_size=chunk_page_size,
            worker_restart_interval=worker_restart_interval,
            workers_per_gpu=workers_per_gpu,
            reserved_cpus=reserved_cpus,
            cpu_workers=cpu_workers,
        )

    def run(self, file_dict: Dict[str, BytesIO]) -> List[Document]:
        """
        Process multiple documents to Markdown format in batch with optimized resource allocation.

        Uses generator-based chunking for immediate task feeding and calculate_optimal_workers
        for dynamic worker allocation based on available GPU/CPU resources.

        Args:
            file_dict: Dictionary mapping filenames (with extensions) to BytesIO file objects

        Returns:
            List of Document objects.
        """
        num_gpus = torch.cuda.device_count()

        # Get total available CPUs
        try:
            total_cpus = len(os.sched_getaffinity(0))
        except Exception:
            total_cpus = os.cpu_count() or 4

        print(f"[DocTool] System Resources: {num_gpus} GPU(s), {total_cpus} CPU(s)")

        # Calculate optimal worker configuration
        num_workers, gpu_ids = calculate_optimal_workers(
            num_gpus=num_gpus,
            total_cpus=total_cpus,
            workers_per_gpu=self.config.workers_per_gpu,
            reserved_cpus=self.config.reserved_cpus
        )

        if gpu_ids:
            print(f"[DocTool] GPU Mode: {num_workers} workers across {num_gpus} GPU(s) ({self.config.workers_per_gpu} workers/GPU)")
            print(f"[DocTool] GPU Assignment: {gpu_ids}")
        else:
            print(f"[DocTool] CPU Mode: {num_workers} workers ({total_cpus - self.config.reserved_cpus} available CPUs)")

        # Initialize document tracker for timing
        doc_tracker = {}

        # Convert config to dict
        config_dict = {
            "do_ocr": self.config.do_ocr,
            "do_table_structure": self.config.do_table_structure,
            "do_formula_enrichment": self.config.do_formula_enrichment,
            "generate_picture_images": self.config.generate_picture_images,
            "images_scale": self.config.images_scale,
            "layout_batch_size": self.config.layout_batch_size,
            "table_batch_size": self.config.table_batch_size,
            "doc_batch_concurrency": self.config.doc_batch_concurrency,
            "chunk_page_size": self.config.chunk_page_size,
            "worker_restart_interval": self.config.worker_restart_interval,
            "cpu_workers": self.config.cpu_workers,
        }

        # Initialize worker manager
        manager = WorkerManager(
            num_workers=num_workers,
            gpu_ids=gpu_ids,
            config_dict=config_dict,
            worker_restart_interval=self.config.worker_restart_interval
        )

        # Start workers immediately - they'll wait for tasks
        manager.start_workers()

        # Feed chunks to queue using generator (eliminates pre-chunking bottleneck)
        # This allows workers to start processing immediately as chunks become available
        print("[DocTool] Starting lazy chunk generation...")

        total_tasks = 0
        chunk_results = []

        # Generator produces chunks on-demand
        chunk_generator = _generate_chunks_lazy(file_dict, self.config.chunk_page_size)

        # Feed tasks to queue as they're generated
        for chunk_data in chunk_generator:
            chunk_filename, chunk_index, original_file_id, chunk_bytes, page_offset, is_pdf, _ = chunk_data

            # Initialize tracker for new files
            if is_pdf and original_file_id not in doc_tracker:
                doc_tracker[original_file_id] = {
                    "total": 0,  # Will be incremented as chunks are generated
                    "received": 0,
                    "start_time": time.perf_counter()
                }

            if is_pdf:
                doc_tracker[original_file_id]["total"] += 1

            # Put task in queue immediately
            task = (chunk_filename, chunk_index, original_file_id, chunk_bytes, page_offset)
            manager.task_queue.put(task)
            total_tasks += 1

        print(f"[DocTool] Generated and queued {total_tasks} total tasks")

        # Collect results
        received_results = 0

        print(f"[DocTool] Waiting for {total_tasks} chunk results...")

        while received_results < total_tasks:
            # Check and restart dead workers
            for i in range(manager.num_workers):
                if not manager.processes[i].is_alive():
                    gpu_id = manager.gpu_ids[i] if manager.gpu_ids else None
                    device_name = f"GPU-{gpu_id}" if gpu_id is not None else f"CPU-Worker{i}"
                    print(f"[Manager] Detected stopped worker {device_name}, restarting...")
                    manager.restart_worker(i, gpu_id)

            try:
                result = manager.result_queue.get(timeout=30)  # Shorter timeout for more frequent worker checks
                
                fid = result.original_file_id
                if fid in doc_tracker:
                    tracker = doc_tracker[fid]
                    tracker["received"] += 1
                    
                    # Calculate and print the time when all chunks have arrived
                    if tracker["received"] == tracker["total"]:
                        end_time = time.perf_counter()
                        elapsed_time = end_time - tracker["start_time"]
                        
                        # elapsed_time
                        print(f"[Processed] {fid} completed in {elapsed_time:.2f} seconds")
                
                chunk_results.append(result)
                received_results += 1

            except Empty:
                # On timeout, check worker status and continue
                # Check if any workers are still alive
                alive_workers = sum(1 for p in manager.processes if p.is_alive())
                if alive_workers == 0 and received_results < total_tasks:
                    print(f"[DocTool] ERROR: All workers died with {total_tasks - received_results} tasks remaining!")
                    break

                continue  # Continue waiting for results

        # Shutdown workers
        manager.shutdown()

        print(f"[DocTool] All chunks processed. Merging results...")

        # Merge chunk results into final documents
        # Generator-based approach processes ALL files (PDF and non-PDF) through workers
        final_documents = self._merge_chunk_results(chunk_results)

        print(f"[DocTool] Processing complete. Total documents: {len(final_documents)}")

        return final_documents

    def _merge_chunk_results(self, chunk_results: List[ChunkResult]) -> List[Document]:
        """
        Merge chunk results back into complete documents.

        Args:
            chunk_results: List of ChunkResult objects

        Returns:
            List of complete Document objects
        """
        # Group by original file ID
        file_chunks: Dict[str, List[ChunkResult]] = {}

        for chunk in chunk_results:
            if chunk.original_file_id not in file_chunks:
                file_chunks[chunk.original_file_id] = []
            file_chunks[chunk.original_file_id].append(chunk)

        # Merge each file's chunks
        documents = []

        for file_id, chunks in file_chunks.items():
            # Sort by chunk index
            chunks.sort(key=lambda x: x.chunk_index)

            # Check for errors
            failed_chunks = [c for c in chunks if not c.success]
            if failed_chunks:
                print(f"[Warning] {file_id} has {len(failed_chunks)} failed chunks:")
                for fc in failed_chunks:
                    print(f"  - Chunk {fc.chunk_index}: {fc.error_msg}")

            # Merge text
            merged_text = "\n\n".join([c.text for c in chunks if c.success])

            # Merge images
            all_images = []
            for chunk in chunks:
                if chunk.success and chunk.images:
                    all_images.extend(chunk.images)

            # Create final document
            doc = Document(
                id=file_id,
                text=merged_text,
                images=all_images if all_images else None
            )

            documents.append(doc)

        return documents


# ex
if __name__ == "__main__":
    # Set multiprocessing start method (required for CUDA)
    try:
        multiprocessing.set_start_method('spawn', force=True)
    except RuntimeError:
        pass

    input_folder = Path("/home/shaush/pdfs")
    output_root = Path("/home/shaush/work/parsed-outputs-cpu-core8")
    log_file_path = output_root / "parsing_log.txt"

    output_root.mkdir(parents=True, exist_ok=True)

    file_list = [p.resolve() for p in input_folder.iterdir() if p.is_file()]
    print(f"Found {len(file_list)} files.")

    # Check available GPUs
    num_gpus = torch.cuda.device_count()
    if num_gpus > 0:
        print(f"Detected {num_gpus} GPU(s). Multi-GPU processing will be used.")
        for i in range(num_gpus):
            gpu_name = torch.cuda.get_device_name(i)
            print(f"  GPU {i}: {gpu_name}")
    else:
        print("No GPU detected. Using CPU mode.")

    processor = DocTool(
        chunk_page_size=10,
        worker_restart_interval=100,  # Increased for long-lived workers (reduces model loading overhead)
        workers_per_gpu=2,             # 2 workers per GPU for better GPU utilization
        reserved_cpus=2,               # Reserve 2 CPUs for system/main process
        cpu_workers=4                  # Fallback for CPU-only mode (deprecated, auto-calculated)
    )

    file_dict = {}
    for file_path in file_list:
        with open(file_path, "rb") as f:
            file_dict[file_path.name] = BytesIO(f.read())
            
    start_time = time.perf_counter()
    
    results = processor.run(file_dict)
    
    total_time = time.perf_counter() - start_time
    print(f"Total parsing time: {total_time:.2f} seconds")

    with open(log_file_path, "w", encoding="utf-8") as log_file:
        log_file.write(f"Batch Processing Report (Multi-GPU/CPU)\n")
        log_file.write(f"Total Files: {len(results)}\n")
        log_file.write(f"Total Time: {total_time:.2f}s\n")
        log_file.write(f"Average Time per File: {total_time/len(results):.2f}s\n")
        log_file.write(f"GPUs Used: {num_gpus if num_gpus > 0 else 'CPU only'}\n")
        log_file.write("="*50 + "\n")

        for doc in results:
            filename = doc.id

            md_content = doc.text

            save_name = Path(filename).stem + ".md"
            save_path = output_root / save_name

            # Markdown 파일 저장
            try:
                with open(save_path, "w", encoding="utf-8") as f:
                    f.write(md_content)

                num_images = len(doc.images) if doc.images else 0
                log_msg = f"[Success] {filename} | Images extracted: {num_images}"
                log_file.write(log_msg + "\n")
                log_file.write(f"   - doc.id: {doc.id}\n")
                log_file.write(f"   - doc.text preview: {doc.text[:100]}...\n")

                if doc.images:
                    first_img = doc.images[0]
                    log_file.write(f"   - Sample Image ID: {first_img.id} ({first_img.mime_type})\n")
                    log_file.write(f"   - Image data length: {len(first_img.data)} bytes\n")

            except Exception as e:
                err_msg = f"[Failed] {filename}: {e}"
                print(err_msg)
                log_file.write(err_msg + "\n")

    print(f"Done! Results saved in '{output_root}'")

