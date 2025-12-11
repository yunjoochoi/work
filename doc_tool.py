import shutil
import time
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, Optional, List, Tuple

from docling.document_converter import DocumentConverter, PdfFormatOption
from docling.datamodel.base_models import InputFormat
from docling.datamodel.pipeline_options import PdfPipelineOptions
from docling_core.types.doc.base import ImageRefMode
from docling_core.types.io import DocumentStream
from docling.datamodel.settings import settings

from dataclasses import dataclass

# for _process_pdf_document
from docling_core.transforms.serializer.markdown import MarkdownDocSerializer
from docling_core.types.doc import DocItem, PictureItem, TableItem
import fitz  

from openpyxl import load_workbook
from openpyxl.utils import range_boundaries
from pptx import Presentation
import pandas as pd
import re

from markitdown import MarkItDown
import mammoth

@dataclass
class ParserConfig:
    """
    Configuration dataclass for DoclingParser.

    Attributes:
        regions_dir: Base directory path for storing extracted images
        do_ocr: Whether to perform OCR on images within documents
        do_table_structure: Whether to detect and preserve table structures
        images_scale: Scale factor for image resolution (higher = better quality)
        layout_batch_size: Batch size for layout detection processing
        table_batch_size: Batch size for table structure recognition
        doc_batch_size: Number of documents to process in parallel
        doc_batch_concurrency: Maximum concurrent document processing threads
    """

    regions_dir: str = "media/regions"  # Base directory for extracted images
    do_ocr: bool = False
    do_table_structure: bool = True  # Enable table structure detection
    images_scale: float = 2.0  # Scale factor for generated images
    layout_batch_size: int = 8  # Batch size for layout detection
    table_batch_size: int = 8  # Batch size for table structure recognition

    doc_batch_size: int = 4  # doc batch
    doc_batch_concurrency: int = 4


class DoclingParser:
    """
    Document parser that converts various formats to Markdown using Docling library.

    Features:
    - Page-wise processing for paginated documents (PDF)
    - Image extraction and organization by page
    - Support for both paginated (PDF) and linear (DOCX) documents
    """

    def __init__(self, config: Optional[ParserConfig] = None):
        """
        Initialize the Docling parser with configuration.

        Args:
            config: Parser configuration. If None, uses default ParserConfig
        """

        self.config = config or ParserConfig()

        settings.perf.doc_batch_size = self.config.doc_batch_size
        settings.perf.doc_batch_concurrency = self.config.doc_batch_concurrency
        pipeline_options = self._create_pipeline_options(self.config)

        self.converter = DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options)
            }
        )

        self.regions_dir = Path(self.config.regions_dir)
        self.regions_dir.mkdir(parents=True, exist_ok=True)

    @staticmethod
    def _create_pipeline_options(config: ParserConfig) -> PdfPipelineOptions:
        """
        Create PDF processing pipeline options from parser configuration.

        Args:
            config: Parser configuration containing pipeline settings

        Returns:
            Configured PdfPipelineOptions instance with OCR, table detection, and batch settings
        """
        options = PdfPipelineOptions()

        options.do_ocr = config.do_ocr
        options.do_table_structure = config.do_table_structure
        options.generate_picture_images = True
        options.images_scale = config.images_scale

        if hasattr(options, "layout_batch_size"):
            options.layout_batch_size = config.layout_batch_size

        if hasattr(options, "table_batch_size"):
            options.table_batch_size = config.table_batch_size

        return options

    def parse(
        self,
        display_name: str,
        file_obj: BytesIO,
    ) -> str:
        """
        Parse a document file to Markdown format with image extraction.

        Extracts images to organized folders and generates markdown with image references.
        Automatically detects document type and applies appropriate processing.

        Args:
            display_name: Display name of the file including extension (e.g., "report.pdf")
            file_obj: BytesIO object containing the document file data

        Returns:
            Markdown formatted text with image references and page separators
        """
        start_time = time.time()

        path_obj = Path(display_name)
        filename = path_obj.stem

        # Read file content once and create separate BytesIO objects for each use
        # This prevents "I/O operation on closed file" errors when Docling closes the stream
        file_obj.seek(0)
        file_bytes = file_obj.read()

        # Create document stream for Docling
        doc_stream = DocumentStream(name=display_name, stream=BytesIO(file_bytes))

        # Convert document
        result = self.converter.convert(doc_stream)
        doc = result.document

        # Convert to markdown text using fresh BytesIO object
        markdown_text = self._convert_to_markdown(
            doc=doc,
            display_name=display_name,
            file_obj=BytesIO(file_bytes)
        )

        elapsed_time = time.time() - start_time
        print(f"Parser execution time: {elapsed_time:.2f} seconds")

        return markdown_text

    def _convert_to_markdown(
        self,
        doc,
        display_name: str,
        file_obj: BytesIO
    ) -> str:
        """
        Convert parsed document to Markdown text based on file type.

        Dispatches to format-specific processing methods:
        - PDF: Hybrid processing (Fitz for images/text + Docling for tables)
        - PPTX: Slide-by-slide with chart extraction (file opened once)
        - Excel: Sheet-by-sheet with chart extraction (file opened once)
        - DOCX: Linear document with image extraction
        - Others: Basic Docling processing

        Args:
            doc: Docling document object containing parsed content
            display_name: File name with extension (e.g., "report.pdf")
            file_obj: BytesIO object containing file data

        Returns:
            Complete markdown text with appropriate headers and content
        """
        path_obj = Path(display_name)
        ext = path_obj.suffix.lower()
        file_key = path_obj.stem

        markdown_text = ""

        # 1. PDF: Hybrid processing (Fitz + Docling)
        if ext == '.pdf':
            markdown_text = self._process_pdf_document(doc, file_key, file_obj)

        # 2. PPTX: Chart extraction + slide-by-slide processing
        elif ext in ['.pptx', '.ppt']:
            markdown_text = self._process_pptx_document(doc, file_key, file_obj)

        # 3. Excel: Sheet names + chart extraction in single pass
        elif ext in ['.xlsx', '.xls', '.xlsm']:
            markdown_text = self._process_excel_document(doc, file_key, file_obj)

        # 4. Word: Linear document with image extraction
        elif ext in ['.docx', '.doc']:
            markdown_text = self._process_docx_document(doc, file_key)

        # 5. Others: Basic processing (Fallback)
        else:
            print(f"[Info] Unknown format {ext}, using default processing.")
            markdown_text = self._process_docx_document(doc, file_key)

        return markdown_text.strip()

    def _process_paginated_document(
        self,
        doc,
        page_num: int,
        file_key: str,
    ) -> str:
        """
        Process a single page from a paginated document (e.g., PDF, Excel).

        Saves page images to page-specific folder (page_XXXX/) and generates markdown for that page.

        Args:
            doc: Docling document object
            page_num: Page number to process (1-indexed)
            file_key: Unique identifier for organizing images

        Returns:
            Markdown text for the specified page with image references
        """
        page_folder = self.regions_dir / file_key / f"page_{page_num:04d}"
        page_folder.mkdir(parents=True, exist_ok=True)

        # Save images for this page and update references
        print(f"  [Saving] Images to {page_folder}")
        page_doc = doc._with_pictures_refs(
            image_dir=page_folder,
            page_no=page_num,
            reference_path=None,
        )

        # Export page to markdown
        markdown_output = page_doc.export_to_markdown(
            page_no=page_num, image_mode=ImageRefMode.REFERENCED
        )

        # Check what files were actually saved
        saved_files = list(page_folder.glob("*")) if page_folder.exists() else []
        print(f"  [Saved] {len(saved_files)} files: {[f.name for f in saved_files]}")

        return markdown_output

    def _process_pdf_document(self, doc, file_key: str, file_obj: BytesIO) -> str:
        """
        [Robust Hybrid Logic]
        마커 형식을 [ID_N]으로 변경하고, 삽입 위치를 안전하게 조정하여 인식률을 높임
        """
        # 1. 파일 데이터 로드
        file_obj.seek(0)
        file_bytes = file_obj.read()
        
        pdf_doc = fitz.open(stream=BytesIO(file_bytes), filetype="pdf")
        serializer = MarkdownDocSerializer(doc=doc)
        
        replacements = {} 
        counter = 0
        
        markdown_parts = []

        for page_num in range(1, doc.num_pages() + 1):
            fitz_page = pdf_doc[page_num - 1]
            page_h = fitz_page.rect.height

            # 페이지 헤더 (선택)
            fitz_page.insert_text(
                (50, 50),
                f"--- Page {page_num} ---",
                fontsize=10,
                color=(0.5, 0.5, 0.5),
            )

            # 이미지 폴더 준비
            page_folder = self.regions_dir / file_key / f"page_{page_num:04d}"
            page_folder.mkdir(parents=True, exist_ok=True)

            # Docling이 이미지 저장 + uri 업데이트
            page_doc = doc._with_pictures_refs(
                image_dir=page_folder,
                page_no=page_num,
                reference_path=Path(f"page_{page_num:04d}"),
            )

            # 👉 1차 루프: redaction만 등록
            items_for_page = []
            for item, _ in page_doc.iterate_items():
                if not isinstance(item, (TableItem, PictureItem)):
                    continue
                if not (item.prov and item.prov[0].bbox):
                    continue

                bbox = item.prov[0].bbox
                l, r = bbox.l, bbox.r
                t = page_h - bbox.t
                b = page_h - bbox.b
                rect = fitz.Rect(l, min(t, b), r, max(t, b))

                items_for_page.append((item, rect))

                # 여기는 redaction만 등록
                fitz_page.add_redact_annot(rect)

            # 페이지당 redaction 한 번만 적용
            if items_for_page:
                fitz_page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_REMOVE)

            # 👉 2차 루프: 마커 삽입 + replacements 등록
            for item, rect in items_for_page:
                try:
                    # content_md 만들기
                    if isinstance(item, PictureItem):
                        if item.image and item.image.uri:
                            img_name = Path(str(item.image.uri)).name
                            rel_path = f"page_{page_num:04d}/{img_name}"
                            content_md = f"\n![Image]({rel_path})\n"
                        else:
                            continue
                    else:
                        content_md = serializer.serialize(item=item).text

                    marker_id = f"ID_{counter}"
                    marker_text = f"[{marker_id}]"
                    replacements[marker_id] = content_md

                    # redaction 이후, 더 이상 이 텍스트를 지우는 작업은 없음
                    fitz_page.insert_text(
                        rect.tl + (0, 10),
                        marker_text,
                        fontsize=8,
                        color=(0, 0, 0),
                        render_mode=0,
                    )
                    counter += 1
                except Exception as e:
                    print(f"Item error on page {page_num}: {e}")
                    continue

            fitz_page.clean_contents()


        # ---------------------------------------------------
        # [Step 3] MarkItDown 변환
        # ---------------------------------------------------
        temp_path = None
        md_output = ""
        
        try:
            temp_filename = f".temp_{file_key}_{int(time.time())}.pdf"
            temp_path = Path(temp_filename).resolve()
            
            pdf_doc.save(str(temp_path), garbage=4, clean=True)
            
            md_tool = MarkItDown()
            result = md_tool.convert(str(temp_path))
            
            if result and result.text_content:
                md_output = result.text_content
            else:
                md_output = ""

        except Exception as e:
            print(f"MarkItDown failed: {e}")
            return ""
        finally:
            pdf_doc.close()
            if temp_path and temp_path.exists():
                try: temp_path.unlink()
                except: pass

        # ---------------------------------------------------
        # [Step 4] 정규식 치환 (Robust Regex)
        # ---------------------------------------------------
        
        # [디버깅] MarkItDown이 읽은 원본 텍스트 일부를 출력해서 마커가 어떻게 생겼나 확인하세요!
        print(f"\n[Debug] MD Raw Content Preview (First 500 chars):\n{md_output[:500]}\n")

        # 패턴: [ID_숫자] 형태를 찾음. (대괄호 이스케이프 주의)
        # MarkItDown이 대괄호 주변에 공백을 넣을 수도 있으므로 \s* 추가
        pattern = re.compile(r"\[\s*(ID_\d+)\s*\]")
        
        def replace_match(match):
            key = match.group(1) # ID_0
            if key in replacements:
                # 앞뒤 개행 2번으로 확실한 분리
                return f"\n\n{replacements[key]}\n\n"
            return match.group(0)

        # 1차 치환 실행
        final_output = pattern.sub(replace_match, md_output)
        
        # 놓친 항목 확인 및 복구
        found_keys = set(re.findall(r"\[\s*(ID_\d+)\s*\]", md_output))
        all_keys = set(replacements.keys())
        lost_keys = all_keys - found_keys
        
        lost_items = []
        for key in lost_keys:
            print(f"[Warning] Recovering lost item [{key}]")
            lost_items.append(replacements[key])

        if lost_items:
            final_output += "\n\n--- [Recovered Elements] ---\n\n"
            final_output += "\n\n".join(lost_items)

        markdown_parts.append(final_output)

        return "".join(markdown_parts)

    def _process_docx_document(self, file_path: str, file_key: str) -> str:
        """
        Process DOCX documents using Mammoth.
        
        Extracts images to a local folder and converts content to Markdown.
        Replaces the Docling implementation with a lightweight Mammoth backend.

        Args:
            file_path: Path to the .docx file (needed for Mammoth to read)
            file_key: Unique identifier for organizing images

        Returns:
            Markdown text with image references
        """
        # 1. 이미지 저장 경로 설정
        images_dir = self.regions_dir / file_key / "images"
        images_dir.mkdir(parents=True, exist_ok=True)
        
        print(f"  [Processing] DOCX with Mammoth: {file_path}")
        print(f"  [Saving] Images to {images_dir}")

        img_counter = 0

        # 2. Mammoth용 커스텀 이미지 핸들러 정의
        def image_handler(image):
            nonlocal img_counter
            img_counter += 1
            
            # 이미지 확장자 추출 (기본값 png)
            content_type = image.content_type
            extension = content_type.split("/")[-1] if "/" in content_type else "png"
            
            # 파일명 생성 (image_1.png, image_2.jpg ...)
            filename = f"image_{img_counter}.{extension}"
            
            # 실제 저장될 물리적 경로 (Absolute Path)
            save_path = images_dir / filename
            
            # 마크다운에 들어갈 상대 경로 (Relative Path)
            # 예: "images/image_1.png" (상황에 따라 절대경로를 써야 할 수도 있음)
            # 여기서는 images_dir 이름이 포함된 상대 경로를 반환한다고 가정
            ref_path = f"images/{filename}" 

            # 이미지 바이너리 저장
            try:
                with image.open() as image_stream:
                    with open(save_path, "wb") as f:
                        f.write(image_stream.read())
            except Exception as e:
                print(f"  [Error] Failed to save image {filename}: {e}")
                return {"src": "", "alt": f"Image_Save_Failed_{img_counter}"}

            # Mammoth에 반환할 속성 (Markdown 변환 시 사용됨)
            return {
                "src": ref_path, 
                "alt": f"Image_{img_counter}"
            }

        # 3. 변환 실행
        try:
            with open(file_path, "rb") as docx_file:
                result = mammoth.convert_to_markdown(
                    docx_file, 
                    convert_image=mammoth.images.img_element(image_handler)
                )
                
            markdown_output = result.value
            messages = result.messages  # 경고 메시지 등 확인 가능

            # 4. 결과 리포트
            saved_files = list(images_dir.glob("*")) if images_dir.exists() else []
            print(f"  [Saved] {len(saved_files)} images extracted.")
            if messages:
                print(f"  [Mammoth Info] {len(messages)} messages reported.")

            return markdown_output

        except Exception as e:
            print(f"  [Error] Mammoth conversion failed: {e}")
            return "" # 또는 에러 raise

    # def _process_docx_document(self, doc, file_key: str) -> str:
    #     """
    #     Process linear (non-paginated) documents like DOCX, TXT, etc.

    #     Saves all images to a single images/ folder and generates markdown.
    #     Used as fallback for unknown document formats.

    #     Args:
    #         doc: Docling document object
    #         file_key: Unique identifier for organizing images

    #     Returns:
    #         Markdown text with image references
    #     """
    #     images_dir = self.regions_dir / file_key / "images"
    #     images_dir.mkdir(parents=True, exist_ok=True) #!

    #     # Save images and update references
    #     print(f"  [Saving] Images to {images_dir}") #!
    #     new_doc = doc._with_pictures_refs(
    #         image_dir=images_dir, page_no=None, reference_path=None
    #     )

    #     markdown_output = new_doc.export_to_markdown(image_mode=ImageRefMode.REFERENCED)

    #     # Check what files were actually saved
    #     saved_files = list(images_dir.glob("*")) if images_dir.exists() else [] #!
    #     print(f"  [Saved] {len(saved_files)} files: {[f.name for f in saved_files]}") #!

    #     return markdown_output


    def _process_pptx_document(self, doc, file_key: str, file_obj: BytesIO) -> str:
        """
        Process PowerPoint documents with chart extraction.

        Opens file ONCE to extract all chart data, then processes each slide
        with Docling-generated text and inserts charts at appropriate positions.

        Args:
            doc: Docling document object
            file_key: Unique identifier for organizing images
            file_obj: BytesIO object containing PPTX file

        Returns:
            Complete markdown text with slide headers and chart data
        """
        # Open file ONCE to extract all charts
        charts_by_page = self._extract_charts_from_pptx(file_obj)

        markdown_parts = []
        num_pages = doc.num_pages()

        for page_num in range(1, num_pages + 1):
            # Generate base text from Docling
            page_text = self._process_paginated_document(doc, page_num, file_key)

            # Insert charts if present on this slide
            if page_num in charts_by_page:
                page_text = self._insert_charts_at_position(page_text, charts_by_page[page_num])

            page_header = f"\n\n- Slide {page_num} -\n\n"
            markdown_parts.append(page_header)
            markdown_parts.append(page_text.strip())

        return "".join(markdown_parts)

    def _process_excel_document(self, doc, file_key: str, file_obj: BytesIO) -> str:
        """
        Process Excel documents with sheet names and chart extraction.

        Opens file ONCE to extract both sheet names and chart data,
        then processes each sheet with Docling-generated text.

        Args:
            doc: Docling document object
            file_key: Unique identifier for organizing images
            file_obj: BytesIO object containing Excel file

        Returns:
            Complete markdown text with sheet headers and chart data
        """
        # Open file ONCE to extract sheet names AND charts
        sheet_names, charts_by_page = self._extract_excel_metadata(file_obj)

        markdown_parts = []
        num_pages = doc.num_pages()

        for page_num in range(1, num_pages + 1):
            # Generate base text from Docling
            page_text = self._process_paginated_document(doc, page_num, file_key)

            # Create header with sheet name if available
            if (page_num - 1) < len(sheet_names):
                sheet_name = sheet_names[page_num - 1]
                header = f"\n\n- Sheet: {sheet_name} -\n\n"
            else:
                header = f"\n\n- Sheet {page_num} -\n\n"

            # Insert charts if present on this sheet
            if page_num in charts_by_page:
                page_text = self._insert_charts_at_position(page_text, charts_by_page[page_num])

            markdown_parts.append(header)
            markdown_parts.append(page_text.strip())

        return "".join(markdown_parts)

    # ==================== Chart Insertion and Extraction Methods ====================

    def _find_chart_insert_position(self, text: str, pre_text: str, post_text: str, start_pos: int) -> int:
        """
        Find optimal position to insert a chart based on surrounding text.

        Uses a three-tier search strategy:
        1. Match both pre_text and post_text (most accurate)
        2. Match only pre_text (insert after)
        3. Match only post_text (insert before)

        Args:
            text: Markdown text to search in
            pre_text: Text that appears before the chart
            post_text: Text that appears after the chart
            start_pos: Position to start searching from

        Returns:
            Insert position, or -1 if no match found
        """
        MAX_DISTANCE = 500  # Maximum characters between pre and post text

        # Strategy 1: Match both pre and post text
        if pre_text and post_text:
            pre_idx = text.find(pre_text, start_pos)
            if pre_idx != -1:
                search_start = pre_idx + len(pre_text)
                post_idx = text.find(post_text, search_start)
                if post_idx != -1 and (post_idx - search_start) < MAX_DISTANCE:
                    return search_start

        # Strategy 2: Match only pre_text (insert after)
        if pre_text:
            pre_idx = text.find(pre_text, start_pos)
            if pre_idx != -1:
                return pre_idx + len(pre_text)

        # Strategy 3: Match only post_text (insert before)
        if post_text:
            post_idx = text.find(post_text, start_pos)
            if post_idx != -1:
                return post_idx

        return -1

    def _insert_charts_at_position(self, markdown_text: str, charts: List[Dict[str, str]]) -> str:
        """
        Insert charts at their original positions in the markdown text.

        Uses context text (pre_text and post_text) to locate the best insertion point.
        Falls back to appending at the end if no matching position is found.

        Args:
            markdown_text: Original markdown text from Docling
            charts: List of chart info dicts with pre_text, post_text, title, and table

        Returns:
            Modified markdown with charts inserted at appropriate positions
        """
        result = markdown_text
        last_pos = 0

        for chart in charts:
            pre_text = chart["pre_text"].strip()
            post_text = chart["post_text"].strip()
            chart_markdown = f"\n\n#### 차트: {chart['title']}\n\n{chart['table']}\n\n"

            # Find optimal insertion position
            insert_idx = self._find_chart_insert_position(result, pre_text, post_text, last_pos)

            if insert_idx != -1:
                # Insert chart at found position
                result = result[:insert_idx] + chart_markdown + result[insert_idx:]
                last_pos = insert_idx + len(chart_markdown)
                print(f"# Chart inserted: {chart['title']} (position: {insert_idx})")
            else:
                # Fallback: append to end
                result += chart_markdown
                print(f"# Chart appended: {chart['title']} (no matching position found)")

        return result

    def _get_chart_context(self, shape, all_shapes) -> Tuple[str, str]:
        """
        Find the closest text above and below a chart.

        Args:
            shape: Chart shape object
            all_shapes: List of all shapes in the slide

        Returns:
            Tuple of (text_above, text_below)
        """
        chart_top = shape.top
        chart_bottom = shape.top + shape.height

        text_above = []
        text_below = []

        for other in all_shapes:
            # Skip if it's the chart itself or has no text
            if other == shape or not other.has_text_frame:
                continue

            text_content = other.text_frame.text.strip()
            if not text_content:
                continue

            other_top = other.top
            other_bottom = other.top + other.height

            # Text above the chart
            if other_bottom < chart_top:
                distance = chart_top - other_bottom
                text_above.append((distance, text_content))

            # Text below the chart
            elif other_top > chart_bottom:
                distance = other_top - chart_bottom
                text_below.append((distance, text_content))

        # Sort by distance (closest first)
        text_above.sort(key=lambda x: x[0])
        text_below.sort(key=lambda x: x[0])

        # Return closest text (empty string if none found)
        context_pre = text_above[0][1] if text_above else ""
        context_post = text_below[0][1] if text_below else ""

        return context_pre, context_post

    def _extract_charts_from_pptx(self, file_obj: BytesIO) -> Dict[int, List[Dict[str, str]]]:
        """
        Extract charts from PPTX file and convert to markdown format.

        Args:
            file_obj: BytesIO object containing PPTX file

        Returns:
            Dictionary mapping page numbers to list of chart info dicts with keys:
                - title: Chart title
                - pre_text: Text above the chart
                - post_text: Text below the chart
                - table: Markdown table of chart data
        """
        charts_by_page = {}

        try:
            file_obj.seek(0)
            prs = Presentation(file_obj)

            for slide_idx, slide in enumerate(prs.slides):
                page_num = slide_idx + 1
                all_shapes = list(slide.shapes)

                for shape in slide.shapes:
                    if shape.has_chart:
                        chart = shape.chart

                        # Extract chart title
                        try:
                            title = chart.chart_title.text_frame.text
                        except:
                            title = "no title"

                        # Find context text above and below
                        pre_text, post_text = self._get_chart_context(shape, all_shapes)

                        # Extract chart data and convert to markdown table
                        try:
                            df = pd.DataFrame()
                            plot = chart.plots[0]
                            cats = [c.label for c in plot.categories]

                            # Generate categories if none exist
                            if not cats:
                                cats = [f"Item {i}" for i in range(len(plot.series[0].values))]

                            df.index = cats
                            for ser in plot.series:
                                df[ser.name] = pd.Series(ser.values, index=cats)
                            md_table = df.to_markdown()
                        except Exception as e:
                            md_table = f"(fail to extract data: {str(e)})"

                        # Store chart info
                        chart_info = {
                            "title": title,
                            "pre_text": pre_text,
                            "post_text": post_text,
                            "table": md_table
                        }

                        # Add to page's chart list
                        if page_num not in charts_by_page:
                            charts_by_page[page_num] = []
                        charts_by_page[page_num].append(chart_info)

            return charts_by_page

        except Exception as e:
            print(f"Error extracting charts from PPTX: {str(e)}")
            return {}

    # ==================== Excel Metadata Extraction Methods ====================

    def _extract_excel_metadata(self, file_obj: BytesIO) -> Tuple[List[str], Dict[int, List[Dict[str, str]]]]:
        """
        Extract sheet names AND charts from Excel file in a single pass.

        Opens the Excel file ONCE and extracts all necessary metadata:
        - Sheet names for headers
        - Chart data with position context

        Args:
            file_obj: BytesIO object containing Excel file

        Returns:
            Tuple of (sheet_names, charts_by_page) where:
                - sheet_names: List of sheet names in order
                - charts_by_page: Dictionary mapping page numbers (sheet numbers) to list of chart info dicts
        """
        sheet_names = []
        charts_by_page = {}

        try:
            file_obj.seek(0)
            # data_only=True: Get calculated values instead of formulas
            wb = load_workbook(file_obj, data_only=True)
            sheet_names = wb.sheetnames

            for page_idx, sheet_name in enumerate(sheet_names):
                sheet = wb[sheet_name]
                page_num = page_idx + 1  # 1-based page number

                # Check if sheet has charts
                charts = getattr(sheet, "charts", []) or getattr(sheet, "_charts", [])

                if not charts:
                    continue

                page_charts = []

                for chart in charts:
                    # 1. Extract chart title
                    title = "Untitled Chart"
                    try:
                        if chart.title:
                            # Title entered directly
                            if hasattr(chart.title, 'tx') and chart.title.tx.rich:
                                title = chart.title.tx.rich.p[0].r[0].t
                            # Title referencing a cell
                            elif hasattr(chart.title, 'tx') and chart.title.tx.strRef:
                                ref_vals = self._get_values_from_ref(wb, chart.title.tx.strRef.f)
                                if ref_vals:
                                    title = ref_vals[0]
                    except Exception:
                        pass

                    # 2. Extract context text (Pre/Post Text) based on chart position
                    pre_text, post_text = self._get_chart_context_excel(sheet, chart)

                    # 3. Extract chart data and convert to markdown
                    md_table = self._resolve_excel_chart_data(wb, chart)

                    page_charts.append({
                        "title": title,
                        "pre_text": pre_text,
                        "post_text": post_text,
                        "table": md_table
                    })

                if page_charts:
                    charts_by_page[page_num] = page_charts

            return sheet_names, charts_by_page

        except Exception as e:
            print(f"Error extracting Excel metadata: {str(e)}")
            return [], {}

    def _get_chart_context_excel(self, sheet, chart) -> Tuple[str, str]:
        """
        Extract text above and below a chart in Excel.

        Args:
            sheet: Excel worksheet object
            chart: Chart object

        Returns:
            Tuple of (text_above, text_below)
        """
        pre_text = ""
        post_text = ""
        try:
            # Get anchor information (TwoCellAnchor recommended)
            anchor = chart.anchor
            if hasattr(anchor, '_from'):
                row_start = anchor._from.row
                col_start = anchor._from.col

                # Pre-text: Check row above chart start
                if row_start > 0:
                    cell = sheet.cell(row=row_start, column=col_start + 1)
                    pre_text = str(cell.value) if cell.value else ""

                # Post-text: Check row below chart end
                row_end = anchor.to.row
                cell = sheet.cell(row=row_end + 2, column=col_start + 1)
                post_text = str(cell.value) if cell.value else ""
        except Exception:
            pass  # Skip if position info unavailable or OneCellAnchor

        return pre_text.strip(), post_text.strip()

    def _resolve_excel_chart_data(self, wb, chart) -> str:
        """
        Parse cell references from chart object and convert actual data
        from Pandas DataFrame to Markdown format.

        Args:
            wb: Workbook object
            chart: Chart object

        Returns:
            Markdown table string
        """
        try:
            data_dict = {}
            categories = []

            # --- 1. Get X-axis (category) data ---
            # Usually use first series category reference as common X-axis
            if len(chart.series) > 0:
                cat_ref = None
                try:
                    # Check string reference (strRef) or numeric reference (numRef)
                    if hasattr(chart.series[0], 'cat'):
                        if chart.series[0].cat and chart.series[0].cat.strRef:
                            cat_ref = chart.series[0].cat.strRef.f
                        elif chart.series[0].cat and chart.series[0].cat.numRef:
                            cat_ref = chart.series[0].cat.numRef.f

                    if cat_ref:
                        categories = self._get_values_from_ref(wb, cat_ref)
                except Exception:
                    pass  # Auto-generate if category extraction fails

            # --- 2. Get Y-axis (series values) data ---
            for series in chart.series:
                # Series name
                series_name = "Series"
                if series.title:
                    if hasattr(series.title, 'tx') and series.title.tx.rich:
                        series_name = series.title.tx.rich.p[0].r[0].t
                    elif hasattr(series.title, 'tx') and series.title.tx.strRef:
                        ref_vals = self._get_values_from_ref(wb, series.title.tx.strRef.f)
                        if ref_vals:
                            series_name = str(ref_vals[0])

                # Actual data values
                vals = []
                if series.val:
                    if series.val.numRef:
                        vals = self._get_values_from_ref(wb, series.val.numRef.f)

                data_dict[series_name] = vals

            # --- 3. Convert to Markdown via Pandas ---
            if not data_dict:
                return "(No chart data reference found)"

            # Match data lengths (based on longest data)
            max_len = max(len(v) for v in data_dict.values())

            # Auto-generate categories if missing or length mismatch
            if not categories or len(categories) != max_len:
                categories = [f"Item {i+1}" for i in range(max_len)]

            # Create DataFrame
            df = pd.DataFrame(data_dict, index=categories[:max_len])

            # Convert to Markdown
            return df.to_markdown()

        except Exception as e:
            return f"(Error parsing chart data: {str(e)})"

    def _get_values_from_ref(self, wb, ref_str: str) -> List[Any]:
        """
        Parse reference string like 'Sheet1!$A$1:$A$5' and return actual values.

        Args:
            wb: Workbook object
            ref_str: Cell reference string

        Returns:
            List of cell values
        """
        try:
            if "!" not in ref_str:
                return []

            sheet_part, cell_part = ref_str.split("!")
            # Remove quotes from sheet name ('Sheet 1' -> Sheet 1)
            sheet_name = sheet_part.replace("'", "")

            if sheet_name not in wb.sheetnames:
                return []

            sheet = wb[sheet_name]

            # Parse range (using openpyxl utility)
            min_col, min_row, max_col, max_row = range_boundaries(cell_part)

            values = []

            # Read cell values in range order
            for row in sheet.iter_rows(min_row=min_row, max_row=max_row,
                                       min_col=min_col, max_col=max_col,
                                       values_only=True):
                for cell in row:
                    # Handle None values (empty cells as 0 or empty string)
                    values.append(cell if cell is not None else "")

            return values

        except Exception as e:
            print(f"Error parsing reference ({ref_str}): {e}")
            return []


    

class DocTool:
    """
    High-level document processing tool with simplified interface.

    Wraps DoclingParser to provide easy batch processing of multiple documents.
    Handles various document formats (PDF, DOCX, etc.) and converts them to Markdown
    with automatic image extraction and organization.
    """

    def __init__(
        self,
        regions_dir: Optional[str] = None,
        do_ocr: bool = False,
        do_table_structure: bool = True,
    ):
        """
        Initialize the document processing tool.

        Args:
            regions_dir: Base directory for storing extracted images (default: "media/regions")
            do_ocr: Whether to perform OCR on images within documents
            do_table_structure: Whether to detect and preserve table structures
        """
        config = ParserConfig(
            regions_dir=regions_dir or "media/regions",
            do_ocr=do_ocr,
            do_table_structure=do_table_structure,
        )

        self._parser = DoclingParser(config=config)

    def run(self, file_dict: Dict[str, BytesIO]) -> Dict[str, str]:
        """
        Process multiple documents to Markdown format in batch.

        Args:
            file_dict: Dictionary mapping filenames (with extensions) to BytesIO file objects

        Returns:
            Dictionary mapping filenames to their converted markdown text content
        """
        results: Dict[str, str] = {}
        for display_name, file_obj in file_dict.items():
            markdown_text = self._parser.parse(display_name, file_obj)
            results[display_name] = markdown_text
        return results


# example
if __name__ == "__main__":
    import os
    # 상단에 필요한 import가 되어 있다고 가정 (Path, BytesIO, DocTool 등)

    # 1. 입력 및 출력 경로 설정
    input_folder = Path("/home/shaush/pdfs")
    output_root = Path("/home/shaush/work/parsed-outputs")
    
    # DocTool 초기화 
    # (regions_dir를 output_root로 설정하여 결과 폴더 내부에 이미지 폴더가 생성되도록 함)
    processor = DocTool(regions_dir=str(output_root))

    # 파일 목록 가져오기
    file_list = [p.resolve() for p in input_folder.iterdir() if p.is_file()]
    print(f"Found {len(file_list)} files.")

    for file_path in file_list:
        file_stem = file_path.stem  # 확장자 뺀 파일명 (예: 'report')
        file_name = file_path.name  # 확장자 포함 파일명 (예: 'report.pdf')
        
        print(f"\n--- Processing: {file_name} ---")

        try:
            # 파일을 메모리에 로드
            with open(file_path, "rb") as f:
                file_bytes = BytesIO(f.read())

            # 변환 실행 (단일 파일)
            results = processor.run({file_name: file_bytes})
            markdown_content = results[file_name]

            # [저장 로직] 
            # 출력 폴더 생성: output_root / 파일명 (확장자 제외)
            target_dir = output_root / file_stem
            target_dir.mkdir(parents=True, exist_ok=True)

            # MD 파일 저장: output_root / 파일명 / 파일명.md
            save_path = target_dir / f"{file_stem}.md"

            with open(save_path, "w", encoding="utf-8") as f:
                f.write(markdown_content)
            
            print(f"Saved to: {save_path}")

        except Exception as e:
            print(f"[Error] Failed to process {file_name}: {e}")
            import traceback
            traceback.print_exc()
            continue
