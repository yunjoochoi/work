import shutil
import time
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, Optional, List, Tuple

from docling.document_converter import DocumentConverter, PdfFormatOption
from docling.datamodel.base_models import InputFormat
from docling.datamodel.pipeline_options import PdfPipelineOptions
from docling_core.types.doc.base import ImageRefMode
from docling_core.types.io import DocumentStream
from docling.datamodel.settings import settings

from dataclasses import dataclass

# PPTX chart extraction dependencies
try:
    from pptx import Presentation
    import pandas as pd
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False

# Excel chart extraction dependencies
try:
    from openpyxl import load_workbook
    from openpyxl.utils.cell import range_boundaries
    EXCEL_SUPPORT = True
except ImportError:
    EXCEL_SUPPORT = False

@dataclass
class ParserConfig:
    """Docling Parser Config"""
    output_base_dir: str = "parsed_output"
    do_ocr: bool = False
    do_table_structure: bool = True   # Enable table structure detection
    images_scale: float = 2.0         # Scale factor for generated images
    layout_batch_size: int = 8        # Batch size for layout detection
    table_batch_size: int = 8         # Batch size for table structure recognition

    doc_batch_size: int = 4 # doc batch
    doc_batch_concurrency: int = 4

class DoclingParser:
    """
    Parse documents to Markdown using Docling. with page-wise processing and image extraction.
    Provides configurable OCR, table structure detection, and image generation.
    """
    def __init__(self, config: ParserConfig):
        """
        Initialize the Docling parser.
        """
        self.config = config
        
        settings.perf.doc_batch_size = self.config.doc_batch_size
        settings.perf.doc_batch_concurrency = self.config.doc_batch_concurrency
        pipeline_options = self._create_pipeline_options(self.config)

        self.converter = DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options)
            }
        )

        self.output_base_dir = Path(self.config.output_base_dir)
        self.output_base_dir.mkdir(exist_ok=True)


    @staticmethod
    def _create_pipeline_options(config: ParserConfig) -> PdfPipelineOptions:
        """Create and configure PDF pipeline options from ParserConfig."""
        options = PdfPipelineOptions()
        
        options.do_ocr = config.do_ocr
        options.do_table_structure = config.do_table_structure
        options.generate_picture_images = True
        options.images_scale = config.images_scale
        
        if hasattr(options, 'layout_batch_size'):
            options.layout_batch_size = config.layout_batch_size
        
        if hasattr(options, 'table_batch_size'):
            options.table_batch_size = config.table_batch_size
            
        return options

    def _input_streams(
        self, 
        file_dict: Dict[str, BytesIO]
    ) -> Tuple[List[DocumentStream], Dict[str, bytes]]:
        """
        file_dict -> DocumentStream
        """
        doc_streams = []
        raw_bytes_map = {}

        for filename, stream in file_dict.items():
            # 1. 스트림 커서 초기화 (안전장치)
            stream.seek(0)
            
            # 2. 바이트 데이터 읽기 (복제)
            # 스트림을 Docling과 Fitz가 공유하면 'Closed file' 에러가 날 수 있으므로
            # 아예 메모리에 bytes로 떠두는 것이 가장 안전합니다.
            file_bytes = stream.read()
            
            # 3. Raw 데이터 저장 (나중에 Fitz/Zipfile이 씀)
            raw_bytes_map[filename] = file_bytes
            
            # 4. Docling용 DocumentStream 생성 (새로운 BytesIO로 감싸서 전달)
            doc_streams.append(
                DocumentStream(name=filename, stream=BytesIO(file_bytes))
            )
            
        return doc_streams, raw_bytes_map

    def parse(
        self,
        file_dict: Dict[str, BytesIO],
        output_dir: Optional[Path] = None,
    ) -> Dict[str, Any]:
        """
        Parse a document to Markdown format.

        Args:
            display_name: Display name of the file including ext
            file_obj: BytesIO object containing the file data
            output_dir: Optional output directory override

        Returns:
            Dictionary containing parsing results with keys:
                - format: File extension
                - markdown: Path to generated markdown file
                - output_dir: Output directory path
                - num_pages: Number of pages in the document
        """
        # TODO:
        start_time = time.time()
        base_dir = Path(output_dir) if output_dir else self.output_base_dir

        # 1. 데이터 준비 (스트림 리스트 + 원본 바이트 맵)
        doc_streams, raw_bytes_map = self._input_streams(file_dict)

        results_map = {}

        # 2. 일괄 변환 (병렬 처리 수행됨)
        # raises_on_error=False로 설정하여 개별 파일 실패가 전체를 멈추지 않도록 함
        conv_results_iter = self.converter.convert_all(doc_streams, raises_on_error=False)

        # 3. 결과 순회 및 후처리
        for result in conv_results_iter:
            filename = result.input.file.name
            
            if result.status.name != "SUCCESS":
                print(f"[실패] {filename}: 변환 중 오류 발생")
                continue

            try:
                # 해당 파일의 원본 바이트 가져오기
                file_bytes = raw_bytes_map.get(filename)
                
                # BytesIO 객체 재생성 (PPTX 추출용)
                file_obj_for_pptx = BytesIO(file_bytes) if file_bytes else None

                # 폴더 생성
                path_obj = Path(filename)
                pure_filename = path_obj.stem
                file_output_dir = base_dir / pure_filename
                file_output_dir.mkdir(parents=True, exist_ok=True)
                
                final_markdown_path = file_output_dir / f"{pure_filename}.md"

                # 마크다운 저장 및 차트 추출 로직 실행
                self._save_markdown_with_page_separation(
                    doc=result.document,
                    final_md_path=final_markdown_path,
                    output_root_dir=file_output_dir,
                    display_name=filename,
                    file_obj=file_obj_for_pptx,
                )

                # 빈 폴더 정리
                self._cleanup_empty_dirs(file_output_dir)
                
                # 결과 딕셔너리에 추가 (파일명: 마크다운 경로)
                results_map[filename] = str(final_markdown_path)
                print(f"[완료] {filename}")

            except Exception as e:
                print(f"[에러] {filename} 후처리 중 실패: {e}")
                traceback.print_exc()

        elapsed_time = time.time() - start_time
        print(f"총 처리 시간: {elapsed_time:.2f} 초")

        return results_map

    def _save_markdown_with_page_separation(
        self, doc, final_md_path: Path, output_root_dir: Path, display_name: str, file_obj: BytesIO
    ) -> None:
        """
        Save document as Markdown with page-wise image organization.

        For paginated documents (PDF, etc), saves images in separate folders per page.
        For linear documents (DOCX), saves all images in a single folder.

        Args:
            doc: Docling document object
            final_md_path: Path to the output markdown file
            output_root_dir: Root directory for output files
            display_name: Original filename for format detection
            file_obj: File object for chart extraction
        """

        num_pages = doc.num_pages()

        # Extract charts if this is a PPTX or Excel file
        charts_by_page = {}
        if self._is_pptx_file(display_name):
            charts_by_page = self._extract_charts_from_pptx(file_obj)
            if charts_by_page:
                print(f"Extracted charts from {len(charts_by_page)} PPTX pages")
        elif self._is_excel_file(display_name):
            charts_by_page = self._extract_charts_from_excel(file_obj)
            if charts_by_page:
                print(f"Extracted charts from {len(charts_by_page)} Excel sheets")

        with open(final_md_path, "w", encoding="utf-8") as file_writer:
            if num_pages == 0:
                self._process_linear_document(doc, file_writer, output_root_dir, final_md_path)
            else:
                self._process_paginated_document(
                    doc, file_writer, output_root_dir, final_md_path, num_pages, charts_by_page
                )


    def _process_linear_document(
        self, doc, file_writer, output_root_dir: Path, final_md_path: Path
    ) -> None:
        """Process documents without pages (DOCX)."""
        images_dir = output_root_dir / "images"

        # Save images and update references
        new_doc = doc._with_pictures_refs(
            image_dir=images_dir, page_no=None, reference_path=final_md_path.parent
        )

        markdown_output = new_doc.export_to_markdown(image_mode=ImageRefMode.REFERENCED)
        file_writer.write(markdown_output)

    def _process_paginated_document(
        self,
        doc,
        file_writer,
        output_root_dir: Path,
        final_md_path: Path,
        num_pages: int,
        charts_by_page: Optional[Dict[int, List[Dict[str, str]]]] = None,
    ) -> None:
        """Process paginated documents (PDF, PPTX, etc)."""
        if charts_by_page is None:
            charts_by_page = {}

        for page_num in range(1, num_pages + 1):
            page_folder = output_root_dir / f"page_{page_num:04d}"

            # Save images for this page and update references
            page_doc = doc._with_pictures_refs(
                image_dir=page_folder,
                page_no=page_num,
                reference_path=final_md_path.parent,
            )

            # Export page to markdown
            markdown_output = page_doc.export_to_markdown(
                page_no=page_num, image_mode=ImageRefMode.REFERENCED
            )

            # Insert charts at their original positions if available
            if page_num in charts_by_page:
                markdown_output = self._insert_charts_at_position(
                    markdown_output, charts_by_page[page_num]
                )

            # Write page header and content
            page_header = f"\n\n- Page {page_num} -\n\n"
            file_writer.write(page_header)
            file_writer.write(markdown_output.strip())

    def _cleanup_empty_dirs(self, output_dir: Path) -> None:
        """Remove empty image directories."""
        if not output_dir.exists():
            return

        for dir_path in output_dir.iterdir():
            if not dir_path.is_dir():
                continue

            if dir_path.name == "images" or dir_path.name.startswith(
                "page_"
            ):
                has_files = any(item.is_file() for item in dir_path.rglob("*"))
                if not has_files:
                    shutil.rmtree(dir_path, ignore_errors=True)

    # ==================== PPTX Chart Extraction Methods ====================

    @staticmethod
    def _is_pptx_file(display_name: str) -> bool:
        """Check if file is a PowerPoint file."""
        return Path(display_name).suffix.lower() in ['.pptx', '.ppt']

    @staticmethod
    def _is_excel_file(display_name: str) -> bool:
        """Check if file is an Excel file."""
        return Path(display_name).suffix.lower() in ['.xlsx', '.xls']

    def _insert_charts_at_position(self, markdown_text: str, charts: List[Dict[str, str]]) -> str:
        """
        Insert charts at their original positions in the markdown text.

        Args:
            markdown_text: Original markdown text from Docling
            charts: List of chart info dicts with pre_text, post_text, title, and table

        Returns:
            Modified markdown with charts inserted at appropriate positions
        """
        result = markdown_text
        last_pos = 0 

        for chart in charts:
            pre_text = chart["pre_text"].strip()
            post_text = chart["post_text"].strip()
            
            chart_markdown = f"\n\n#### 차트: {chart['title']}\n\n{chart['table']}\n\n"
            
            inserted = False
            insert_idx = -1
            
            # Case 1: When both pre_text and post_text are present
            if pre_text and post_text:
                curr_pre_idx = result.find(pre_text, last_pos)
                
                if curr_pre_idx != -1:
                    search_start = curr_pre_idx + len(pre_text)
                    curr_post_idx = result.find(post_text, search_start)
                    
                    if curr_post_idx != -1 and (curr_post_idx - search_start) < 500:
                        insert_idx = search_start
                        inserted = True
            
            # Case 2: Search using only Pre_text (if Case 1 fails)
            if not inserted and pre_text:
                curr_pre_idx = result.find(pre_text, last_pos)
                if curr_pre_idx != -1:
                    insert_idx = curr_pre_idx + len(pre_text)
                    inserted = True

            # Case 3: Search using only Post_text (if Case 2 fails)
            if not inserted and post_text:
                curr_post_idx = result.find(post_text, last_pos)
                if curr_post_idx != -1:
                    insert_idx = curr_post_idx # insert in front of post_text 
                    inserted = True

            # insert chart
            if inserted and insert_idx != -1:
                result = result[:insert_idx] + chart_markdown + result[insert_idx:]
                last_pos = insert_idx + len(chart_markdown) # update cursor
                
                print(f"# chart insert success: {chart['title']} (pos: {insert_idx})")
            
            else:
                # if matching fails -> append last
                result += f"{chart_markdown}"

        return result

    def _get_chart_context(self, shape, all_shapes) -> Tuple[str, str]:
        """
        Find the closest text above and below a chart.

        Args:
            shape: Chart shape object
            all_shapes: List of all shapes in the slide

        Returns:
            Tuple of (text_above, text_below)
        """
        chart_top = shape.top
        chart_bottom = shape.top + shape.height

        text_above = []
        text_below = []

        for other in all_shapes:
            # Skip if it's the chart itself or has no text
            if other == shape or not other.has_text_frame:
                continue

            text_content = other.text_frame.text.strip()
            if not text_content:
                continue

            other_top = other.top
            other_bottom = other.top + other.height

            # Text above the chart
            if other_bottom < chart_top:
                distance = chart_top - other_bottom
                text_above.append((distance, text_content))

            # Text below the chart
            elif other_top > chart_bottom:
                distance = other_top - chart_bottom
                text_below.append((distance, text_content))

        # Sort by distance (closest first)
        text_above.sort(key=lambda x: x[0])
        text_below.sort(key=lambda x: x[0])

        # Return closest text (empty string if none found)
        context_pre = text_above[0][1] if text_above else ""
        context_post = text_below[0][1] if text_below else ""

        return context_pre, context_post

    def _extract_charts_from_pptx(self, file_obj: BytesIO) -> Dict[int, List[Dict[str, str]]]:
        """
        Extract charts from PPTX file and convert to markdown format.

        Args:
            file_obj: BytesIO object containing PPTX file

        Returns:
            Dictionary mapping page numbers to list of chart info dicts with keys:
                - title: Chart title
                - pre_text: Text above the chart
                - post_text: Text below the chart
                - table: Markdown table of chart data
        """
        if not PPTX_SUPPORT:
            print("Warning: python-pptx or pandas not installed. Chart extraction disabled.")
            return {}

        charts_by_page = {}

        try:
            file_obj.seek(0)
            prs = Presentation(file_obj)

            for slide_idx, slide in enumerate(prs.slides):
                page_num = slide_idx + 1
                all_shapes = list(slide.shapes)

                for shape in slide.shapes:
                    if shape.has_chart:
                        chart = shape.chart

                        # Extract chart title
                        try:
                            title = chart.chart_title.text_frame.text
                        except:
                            title = "no title"

                        # Find context text above and below
                        pre_text, post_text = self._get_chart_context(shape, all_shapes)

                        # Extract chart data and convert to markdown table
                        try:
                            df = pd.DataFrame()
                            plot = chart.plots[0]
                            cats = [c.label for c in plot.categories]

                            # Generate categories if none exist
                            if not cats:
                                cats = [f"Item {i}" for i in range(len(plot.series[0].values))]

                            df.index = cats
                            for ser in plot.series:
                                df[ser.name] = pd.Series(ser.values, index=cats)
                            md_table = df.to_markdown()
                        except Exception as e:
                            md_table = f"(fail to extract data: {str(e)})"

                        # Store chart info
                        chart_info = {
                            "title": title,
                            "pre_text": pre_text,
                            "post_text": post_text,
                            "table": md_table
                        }

                        # Add to page's chart list
                        if page_num not in charts_by_page:
                            charts_by_page[page_num] = []
                        charts_by_page[page_num].append(chart_info)

            return charts_by_page

        except Exception as e:
            print(f"Error extracting charts from PPTX: {str(e)}")
            return {}

    # ==================== Excel Chart Extraction Methods ====================

    def _extract_charts_from_excel(self, file_obj: BytesIO) -> Dict[int, List[Dict[str, str]]]:
        """
        Extract charts from Excel file and convert to markdown format.

        Args:
            file_obj: BytesIO object containing Excel file

        Returns:
            Dictionary mapping page numbers (sheet numbers) to list of chart info dicts with keys:
                - title: Chart title
                - pre_text: Text above the chart
                - post_text: Text below the chart
                - table: Markdown table of chart data
        """
        if not EXCEL_SUPPORT:
            print("Warning: openpyxl not installed. Excel chart extraction disabled.")
            return {}

        charts_by_page = {}

        try:
            file_obj.seek(0)
            # data_only=True: Get calculated values instead of formulas
            wb = load_workbook(file_obj, data_only=True)

            for page_idx, sheet_name in enumerate(wb.sheetnames):
                sheet = wb[sheet_name]
                page_num = page_idx + 1  # 1-based page number

                # Check if sheet has charts
                charts = getattr(sheet, "charts", []) or getattr(sheet, "_charts", [])

                if not charts:
                    continue

                page_charts = []

                for chart in charts:
                    # 1. Extract chart title
                    title = "Untitled Chart"
                    try:
                        if chart.title:
                            # Title entered directly
                            if hasattr(chart.title, 'tx') and chart.title.tx.rich:
                                title = chart.title.tx.rich.p[0].r[0].t
                            # Title referencing a cell
                            elif hasattr(chart.title, 'tx') and chart.title.tx.strRef:
                                ref_vals = self._get_values_from_ref(wb, chart.title.tx.strRef.f)
                                if ref_vals:
                                    title = ref_vals[0]
                    except Exception:
                        pass

                    # 2. Extract context text (Pre/Post Text) based on chart position
                    pre_text, post_text = self._get_chart_context_excel(sheet, chart)

                    # 3. Extract chart data and convert to markdown
                    md_table = self._resolve_excel_chart_data(wb, chart)

                    page_charts.append({
                        "title": title,
                        "pre_text": pre_text,
                        "post_text": post_text,
                        "table": md_table
                    })

                if page_charts:
                    charts_by_page[page_num] = page_charts

            return charts_by_page

        except Exception as e:
            print(f"Error extracting charts from Excel: {str(e)}")
            return {}

    def _get_chart_context_excel(self, sheet, chart) -> Tuple[str, str]:
        """
        Extract text above and below a chart in Excel.

        Args:
            sheet: Excel worksheet object
            chart: Chart object

        Returns:
            Tuple of (text_above, text_below)
        """
        pre_text = ""
        post_text = ""
        try:
            # Get anchor information (TwoCellAnchor recommended)
            anchor = chart.anchor
            if hasattr(anchor, '_from'):
                row_start = anchor._from.row
                col_start = anchor._from.col

                # Pre-text: Check row above chart start
                if row_start > 0:
                    cell = sheet.cell(row=row_start, column=col_start + 1)
                    pre_text = str(cell.value) if cell.value else ""

                # Post-text: Check row below chart end
                row_end = anchor.to.row
                cell = sheet.cell(row=row_end + 2, column=col_start + 1)
                post_text = str(cell.value) if cell.value else ""
        except Exception:
            pass  # Skip if position info unavailable or OneCellAnchor

        return pre_text.strip(), post_text.strip()

    def _resolve_excel_chart_data(self, wb, chart) -> str:
        """
        Parse cell references from chart object and convert actual data
        from Pandas DataFrame to Markdown format.

        Args:
            wb: Workbook object
            chart: Chart object

        Returns:
            Markdown table string
        """
        try:
            data_dict = {}
            categories = []

            # --- 1. Get X-axis (category) data ---
            # Usually use first series category reference as common X-axis
            if len(chart.series) > 0:
                cat_ref = None
                try:
                    # Check string reference (strRef) or numeric reference (numRef)
                    if hasattr(chart.series[0], 'cat'):
                        if chart.series[0].cat and chart.series[0].cat.strRef:
                            cat_ref = chart.series[0].cat.strRef.f
                        elif chart.series[0].cat and chart.series[0].cat.numRef:
                            cat_ref = chart.series[0].cat.numRef.f

                    if cat_ref:
                        categories = self._get_values_from_ref(wb, cat_ref)
                except Exception:
                    pass  # Auto-generate if category extraction fails

            # --- 2. Get Y-axis (series values) data ---
            for series in chart.series:
                # Series name
                series_name = "Series"
                if series.title:
                    if hasattr(series.title, 'tx') and series.title.tx.rich:
                        series_name = series.title.tx.rich.p[0].r[0].t
                    elif hasattr(series.title, 'tx') and series.title.tx.strRef:
                        ref_vals = self._get_values_from_ref(wb, series.title.tx.strRef.f)
                        if ref_vals:
                            series_name = str(ref_vals[0])

                # Actual data values
                vals = []
                if series.val:
                    if series.val.numRef:
                        vals = self._get_values_from_ref(wb, series.val.numRef.f)

                data_dict[series_name] = vals

            # --- 3. Convert to Markdown via Pandas ---
            if not data_dict:
                return "(차트 데이터 참조를 찾을 수 없음)"

            # Match data lengths (based on longest data)
            max_len = max(len(v) for v in data_dict.values())

            # Auto-generate categories if missing or length mismatch
            if not categories or len(categories) != max_len:
                categories = [f"Item {i+1}" for i in range(max_len)]

            # Create DataFrame
            df = pd.DataFrame(data_dict, index=categories[:max_len])

            # Convert to Markdown
            return df.to_markdown()

        except Exception as e:
            return f"(데이터 파싱 중 오류 발생: {str(e)})"

    def _get_values_from_ref(self, wb, ref_str: str) -> List[Any]:
        """
        Parse reference string like 'Sheet1!$A$1:$A$5' and return actual values.

        Args:
            wb: Workbook object
            ref_str: Cell reference string

        Returns:
            List of cell values
        """
        try:
            if "!" not in ref_str:
                return []

            sheet_part, cell_part = ref_str.split("!")
            # Remove quotes from sheet name ('Sheet 1' -> Sheet 1)
            sheet_name = sheet_part.replace("'", "")

            if sheet_name not in wb.sheetnames:
                return []

            sheet = wb[sheet_name]

            # Parse range (using openpyxl utility)
            min_col, min_row, max_col, max_row = range_boundaries(cell_part)

            values = []

            # Read cell values in range order
            for row in sheet.iter_rows(min_row=min_row, max_row=max_row,
                                       min_col=min_col, max_col=max_col,
                                       values_only=True):
                for cell in row:
                    # Handle None values (empty cells as 0 or empty string)
                    values.append(cell if cell is not None else "")

            return values

        except Exception as e:
            print(f"참조 해석 오류 ({ref_str}): {e}")
            return []

class DocumentProcessor:
    """
    High-level document processor using Docling parser.

    Provides a simple interface for converting documents to Markdown format.
    Supports batch processing of multiple documents.
    """

    def __init__(
        self,
        output_base_dir: Optional[str] = None,
        do_ocr: bool = False,
        do_table_structure: bool = True,
    ):
        """
        Initialize the document processor.

        Args:
            output_base_dir: Base directory for output files
            do_ocr: Enable OCR processing
            do_table_structure: Enable table structure detection
        """
        config = ParserConfig(
            output_base_dir=output_base_dir or "parsed_output",
            do_ocr=do_ocr,
            do_table_structure=do_table_structure
        )
        
        self._parser = DoclingParser(config=config)

    

    def process_file(self, display_name: str, file_obj: BytesIO) -> Dict[str, Any]:
        """
        Process a single document to Markdown.

        Args:
            display_name: Filename with extension
            file_obj: File object containing document data

        Returns:
            Dictionary with parsing results containing:
                - format: File extension
                - markdown: Path to generated markdown file
                - output_dir: Output directory path
                - num_pages: Number of pages
        """
        return self._parser.parse(display_name, file_obj)

    def process(self, file_dict: Dict[str, BytesIO]) -> Dict[str, str]:
        """
        Process multiple documents to Markdown.

        Args:
            file_dict: Dictionary mapping filenames to BytesIO objects

        Returns:
            Dictionary mapping filenames to markdown file paths
        """
        return self._parser.parse(file_dict)



# example
if __name__ == "__main__":
    folder = Path("/data/yunju/pdfs")

    file_list = [p.resolve() for p in folder.iterdir() if p.is_file()]
    print("Found files:", file_list)

    processor = DocumentProcessor()

    file_dict = {}

    for file_path in file_list:
        with open(file_path, "rb") as f:
                file_dict[file_path.name] = BytesIO(f.read())

    results = processor.process(file_dict)

