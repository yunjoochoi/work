import shutil
import time
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, Optional, List, Tuple

from docling.document_converter import DocumentConverter, PdfFormatOption
from docling.datamodel.base_models import InputFormat
from docling.datamodel.pipeline_options import PdfPipelineOptions
from docling_core.types.doc.base import ImageRefMode
from docling_core.types.io import DocumentStream
from docling.datamodel.settings import settings

from dataclasses import dataclass

# PPTX chart extraction dependencies
try:
    from pptx import Presentation
    import pandas as pd
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False

@dataclass
class ParserConfig:
    """Docling Parser Config"""
    output_base_dir: str = "parsed_output"
    do_ocr: bool = False
    do_table_structure: bool = True   # Enable table structure detection
    images_scale: float = 2.0         # Scale factor for generated images
    layout_batch_size: int = 8        # Batch size for layout detection
    table_batch_size: int = 8         # Batch size for table structure recognition

    doc_batch_size: int = 4 # doc batch
    doc_batch_concurrency: int = 4

class DoclingParser:
    """
    Parse documents to Markdown using Docling. with page-wise processing and image extraction.
    Provides configurable OCR, table structure detection, and image generation.
    """
    def __init__(self, config: Optional[ParserConfig] = None):
        """
        Initialize the Docling parser.
        """

        self.config = config or ParserConfig()
        
        settings.perf.doc_batch_size = self.config.doc_batch_size
        settings.perf.doc_batch_concurrency = self.config.doc_batch_concurrency
        pipeline_options = self._create_pipeline_options(self.config)

        self.converter = DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options)
            }
        )

        self.output_base_dir = Path(self.config.output_base_dir)
        self.output_base_dir.mkdir(exist_ok=True)


    @staticmethod
    def _create_pipeline_options(config: ParserConfig) -> PdfPipelineOptions:
        """Create and configure PDF pipeline options from ParserConfig."""
        options = PdfPipelineOptions()
        
        options.do_ocr = config.do_ocr
        options.do_table_structure = config.do_table_structure
        options.generate_picture_images = True
        options.images_scale = config.images_scale
        
        if hasattr(options, 'layout_batch_size'):
            options.layout_batch_size = config.layout_batch_size
        
        if hasattr(options, 'table_batch_size'):
            options.table_batch_size = config.table_batch_size
            
        return options

    def parse(
        self,
        display_name: str,
        file_obj: BytesIO,
        output_dir: Optional[Path] = None,
    ) -> Dict[str, Any]:
        """
        Parse a document to Markdown format.

        Args:
            display_name: Display name of the file including ext
            file_obj: BytesIO object containing the file data
            output_dir: Optional output directory override

        Returns:
            Dictionary containing parsing results with keys:
                - format: File extension
                - markdown: Path to generated markdown file
                - output_dir: Output directory path
                - num_pages: Number of pages in the document
        """
        # TODO:
        start_time = time.time()

        path_obj = Path(display_name)
        filename = path_obj.stem

        base_dir = Path(output_dir) if output_dir else self.output_base_dir
        file_output_dir = base_dir / filename
        file_output_dir.mkdir(parents=True, exist_ok=True)

        # Create document stream for Docling
        doc_stream = DocumentStream(name=display_name, stream=file_obj)

        # Convert document
        result = self.converter.convert(doc_stream)
        doc = result.document

        # Save markdown with page separation
        final_markdown_path = file_output_dir / f"{filename}.md"
        self._save_markdown_with_page_separation(
            doc=doc,
            final_md_path=final_markdown_path,
            output_root_dir=file_output_dir,
            display_name=display_name,
            file_obj=file_obj,
        )

        # Clean up empty directories
        self._cleanup_empty_dirs(file_output_dir)

        # TODO:
        elapsed_time = time.time() - start_time
        print(f"Parser execution time: {elapsed_time:.2f} seconds")

        return {
            "format": path_obj.suffix,
            "markdown": str(final_markdown_path),
            "output_dir": str(file_output_dir),
            "num_pages": doc.num_pages(),
        }

    def _save_markdown_with_page_separation(
        self, doc, final_md_path: Path, output_root_dir: Path, display_name: str, file_obj: BytesIO
    ) -> None:
        """
        Save document as Markdown with page-wise image organization.

        For paginated documents (PDF, etc), saves images in separate folders per page.
        For linear documents (DOCX), saves all images in a single folder.

        Args:
            doc: Docling document object
            final_md_path: Path to the output markdown file
            output_root_dir: Root directory for output files
            display_name: Original filename for format detection
            file_obj: File object for chart extraction
        """

        num_pages = doc.num_pages()

        # Extract charts if this is a PPTX file
        charts_by_page = {}
        if self._is_pptx_file(display_name):
            charts_by_page = self._extract_charts_from_pptx(file_obj)
            if charts_by_page:
                print(f"Extracted charts from {len(charts_by_page)} pages")

        with open(final_md_path, "w", encoding="utf-8") as file_writer:
            if num_pages == 0:
                self._process_linear_document(doc, file_writer, output_root_dir, final_md_path)
            else:
                self._process_paginated_document(
                    doc, file_writer, output_root_dir, final_md_path, num_pages, charts_by_page
                )
        


    def _process_linear_document(
        self, doc, file_writer, output_root_dir: Path, final_md_path: Path
    ) -> None:
        """Process documents without pages (DOCX)."""
        images_dir = output_root_dir / "images"

        # Save images and update references
        new_doc = doc._with_pictures_refs(
            image_dir=images_dir, page_no=None, reference_path=final_md_path.parent
        )

        markdown_output = new_doc.export_to_markdown(image_mode=ImageRefMode.REFERENCED)
        file_writer.write(markdown_output)

    def _process_paginated_document(
        self,
        doc,
        file_writer,
        output_root_dir: Path,
        final_md_path: Path,
        num_pages: int,
        charts_by_page: Optional[Dict[int, List[Dict[str, str]]]] = None,
    ) -> None:
        """Process paginated documents (PDF, PPTX, etc)."""
        if charts_by_page is None:
            charts_by_page = {}

        for page_num in range(1, num_pages + 1):
            page_folder = output_root_dir / f"page_{page_num:04d}"

            # Save images for this page and update references
            page_doc = doc._with_pictures_refs(
                image_dir=page_folder,
                page_no=page_num,
                reference_path=final_md_path.parent,
            )

            # Export page to markdown
            markdown_output = page_doc.export_to_markdown(
                page_no=page_num, image_mode=ImageRefMode.REFERENCED
            )

            # Insert charts at their original positions if available
            if page_num in charts_by_page:
                markdown_output = self._insert_charts_at_position(
                    markdown_output, charts_by_page[page_num]
                )

            # Write page header and content
            page_header = f"\n\n- Page {page_num} -\n\n"
            file_writer.write(page_header)
            file_writer.write(markdown_output.strip())

    def _cleanup_empty_dirs(self, output_dir: Path) -> None:
        """Remove empty image directories."""
        if not output_dir.exists():
            return

        for dir_path in output_dir.iterdir():
            if not dir_path.is_dir():
                continue

            if dir_path.name == "images" or dir_path.name.startswith(
                "page_"
            ):
                has_files = any(item.is_file() for item in dir_path.rglob("*"))
                if not has_files:
                    shutil.rmtree(dir_path, ignore_errors=True)

    # ==================== PPTX Chart Extraction Methods ====================

    @staticmethod
    def _is_pptx_file(display_name: str) -> bool:
        """Check if file is a PowerPoint file."""
        return Path(display_name).suffix.lower() in ['.pptx', '.ppt']

    def _insert_charts_at_position(self, markdown_text: str, charts: List[Dict[str, str]]) -> str:
        """
        Insert charts at their original positions in the markdown text.

        Args:
            markdown_text: Original markdown text from Docling
            charts: List of chart info dicts with pre_text, post_text, title, and table

        Returns:
            Modified markdown with charts inserted at appropriate positions
        """
        result = markdown_text

        for chart in charts:
            pre_text = chart["pre_text"].strip()
            post_text = chart["post_text"].strip()

            # Create chart markdown
            chart_markdown = f"""

#### 차트: {chart["title"]}

{chart["table"]}

"""

            # Try to find insertion point based on context
            inserted = False

            # Case 1: Both pre and post text exist
            if pre_text and post_text:
                # Look for pattern: pre_text ... post_text
                # Insert chart between them
                pre_idx = result.find(pre_text)
                if pre_idx != -1:
                    # Find post_text after pre_text
                    search_start = pre_idx + len(pre_text)
                    post_idx = result.find(post_text, search_start)
                    if post_idx != -1:
                        # Insert chart between pre and post
                        insert_pos = pre_idx + len(pre_text)
                        result = result[:insert_pos] + chart_markdown + result[insert_pos:]
                        inserted = True

            # Case 2: Only pre_text exists
            if not inserted and pre_text:
                pre_idx = result.find(pre_text)
                if pre_idx != -1:
                    insert_pos = pre_idx + len(pre_text)
                    result = result[:insert_pos] + chart_markdown + result[insert_pos:]
                    inserted = True

            # Case 3: Only post_text exists
            if not inserted and post_text:
                post_idx = result.find(post_text)
                if post_idx != -1:
                    result = result[:post_idx] + chart_markdown + result[post_idx:]
                    inserted = True

            # Case 4: No context - append at the end
            if not inserted:
                result += f"\n\n### 차트 정보\n{chart_markdown}"

        return result

    def _get_chart_context(self, shape, all_shapes) -> Tuple[str, str]:
        """
        Find the closest text above and below a chart.

        Args:
            shape: Chart shape object
            all_shapes: List of all shapes in the slide

        Returns:
            Tuple of (text_above, text_below)
        """
        chart_top = shape.top
        chart_bottom = shape.top + shape.height

        text_above = []
        text_below = []

        for other in all_shapes:
            # Skip if it's the chart itself or has no text
            if other == shape or not other.has_text_frame:
                continue

            text_content = other.text_frame.text.strip()
            if not text_content:
                continue

            other_top = other.top
            other_bottom = other.top + other.height

            # Text above the chart
            if other_bottom < chart_top:
                distance = chart_top - other_bottom
                text_above.append((distance, text_content))

            # Text below the chart
            elif other_top > chart_bottom:
                distance = other_top - chart_bottom
                text_below.append((distance, text_content))

        # Sort by distance (closest first)
        text_above.sort(key=lambda x: x[0])
        text_below.sort(key=lambda x: x[0])

        # Return closest text (empty string if none found)
        context_pre = text_above[0][1] if text_above else ""
        context_post = text_below[0][1] if text_below else ""

        return context_pre, context_post

    def _extract_charts_from_pptx(self, file_obj: BytesIO) -> Dict[int, List[Dict[str, str]]]:
        """
        Extract charts from PPTX file and convert to markdown format.

        Args:
            file_obj: BytesIO object containing PPTX file

        Returns:
            Dictionary mapping page numbers to list of chart info dicts with keys:
                - title: Chart title
                - pre_text: Text above the chart
                - post_text: Text below the chart
                - table: Markdown table of chart data
        """
        if not PPTX_SUPPORT:
            print("Warning: python-pptx or pandas not installed. Chart extraction disabled.")
            return {}

        charts_by_page = {}

        try:
            # Reset file pointer
            file_obj.seek(0)
            prs = Presentation(file_obj)

            for slide_idx, slide in enumerate(prs.slides):
                page_num = slide_idx + 1
                all_shapes = list(slide.shapes)

                for shape in slide.shapes:
                    if shape.has_chart:
                        chart = shape.chart

                        # Extract chart title
                        try:
                            title = chart.chart_title.text_frame.text
                        except:
                            title = "제목 없음"

                        # Find context text above and below
                        pre_text, post_text = self._get_chart_context(shape, all_shapes)

                        # Extract chart data and convert to markdown table
                        try:
                            df = pd.DataFrame()
                            plot = chart.plots[0]
                            cats = [c.label for c in plot.categories]

                            # Generate categories if none exist
                            if not cats:
                                cats = [f"Item {i}" for i in range(len(plot.series[0].values))]

                            df.index = cats
                            for ser in plot.series:
                                df[ser.name] = pd.Series(ser.values, index=cats)
                            md_table = df.to_markdown()
                        except Exception as e:
                            md_table = f"(데이터 추출 실패: {str(e)})"

                        # Store chart info
                        chart_info = {
                            "title": title,
                            "pre_text": pre_text,
                            "post_text": post_text,
                            "table": md_table
                        }

                        # Add to page's chart list
                        if page_num not in charts_by_page:
                            charts_by_page[page_num] = []
                        charts_by_page[page_num].append(chart_info)

            return charts_by_page

        except Exception as e:
            print(f"Error extracting charts from PPTX: {str(e)}")
            return {}

class DocumentProcessor:
    """
    High-level document processor using Docling parser.

    Provides a simple interface for converting documents to Markdown format.
    Supports batch processing of multiple documents.
    """

    def __init__(
        self,
        output_base_dir: Optional[str] = None,
        do_ocr: bool = False,
        do_table_structure: bool = True,
    ):
        """
        Initialize the document processor.

        Args:
            output_base_dir: Base directory for output files
            do_ocr: Enable OCR processing
            do_table_structure: Enable table structure detection
        """
        config = ParserConfig(
            output_base_dir=output_base_dir or "parsed_output",
            do_ocr=do_ocr,
            do_table_structure=do_table_structure
        )
        
        self._parser = DoclingParser(config=config)

    def process_file(self, display_name: str, file_obj: BytesIO) -> Dict[str, Any]:
        """
        Process a single document to Markdown.

        Args:
            display_name: Filename with extension
            file_obj: File object containing document data

        Returns:
            Dictionary with parsing results containing:
                - format: File extension
                - markdown: Path to generated markdown file
                - output_dir: Output directory path
                - num_pages: Number of pages
        """
        return self._parser.parse(display_name, file_obj)

    def process(self, file_dict: Dict[str, BytesIO]) -> Dict[str, str]:
        """
        Process multiple documents to Markdown.

        Args:
            file_dict: Dictionary mapping filenames to BytesIO objects

        Returns:
            Dictionary mapping filenames to markdown file paths
        """
        results: Dict[str, str] = {}
        for display_name, file_obj in file_dict.items():
            parse_result = self._parser.parse(display_name, file_obj)
            results[display_name] = parse_result["markdown"]
        return results



# example
if __name__ == "__main__":
    folder = Path("/data/yunju/pdfs")

    file_list = [p.resolve() for p in folder.iterdir() if p.is_file()]
    print("Found files:", file_list)

    processor = DocumentProcessor()

    file_dict = {}

    for file_path in file_list:
        with open(file_path, "rb") as f:
                file_dict[file_path.name] = BytesIO(f.read())

    print(f"\nProcessing {len(file_dict)} file(s)...")
    start_time = time.time()
    results = processor.process(file_dict)
    elapsed_time = time.time() - start_time
    print(f"총 처리 시간: {elapsed_time:.2f} 초")


